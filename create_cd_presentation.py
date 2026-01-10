#!/usr/bin/env python3
"""
Harness CD Workflow Presentation Generator
Creates a comprehensive PowerPoint presentation for VM and ECS deployments
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Color scheme
PRIMARY_COLOR = RGBColor(0, 102, 204)      # Blue
SECONDARY_COLOR = RGBColor(0, 168, 107)    # Green
ACCENT_COLOR = RGBColor(255, 102, 0)       # Orange
DARK_GRAY = RGBColor(51, 51, 51)           # Dark Gray
LIGHT_GRAY = RGBColor(242, 242, 242)       # Light Gray

def create_presentation():
    """Create the PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    create_title_slide(prs)

    # Slide 2: Agenda
    create_agenda_slide(prs)

    # Slide 3: Introduction to Harness CD
    create_intro_slide(prs)

    # Slide 4: CD Flow Overview
    create_cd_flow_overview_slide(prs)

    # Slide 5: VM Deployment Architecture
    create_vm_architecture_slide(prs)

    # Slide 6: VM Deployment Dataflow
    create_vm_dataflow_slide(prs)

    # Slide 7: VM Deployment Components
    create_vm_components_slide(prs)

    # Slide 8: VM Deployment Pipeline Steps
    create_vm_pipeline_steps_slide(prs)

    # Slide 9: VM Pre-Deployment Phase
    create_vm_pre_deployment_slide(prs)

    # Slide 10: VM Deployment Execution
    create_vm_deployment_execution_slide(prs)

    # Slide 11: VM Post-Deployment Phase
    create_vm_post_deployment_slide(prs)

    # Slide 12: ECS Deployment Architecture
    create_ecs_architecture_slide(prs)

    # Slide 13: ECS Deployment Dataflow
    create_ecs_dataflow_slide(prs)

    # Slide 14: ECS Components
    create_ecs_components_slide(prs)

    # Slide 15: ECS Task Definition
    create_ecs_task_definition_slide(prs)

    # Slide 16: ECS Deployment Strategies
    create_ecs_deployment_strategies_slide(prs)

    # Slide 17: ECS Rolling Deployment
    create_ecs_rolling_deployment_slide(prs)

    # Slide 18: ECS Blue-Green Deployment
    create_ecs_blue_green_slide(prs)

    # Slide 19: ECS Canary Deployment
    create_ecs_canary_deployment_slide(prs)

    # Slide 20: Deployment Comparison
    create_deployment_comparison_slide(prs)

    # Slide 21: Best Practices - VM
    create_vm_best_practices_slide(prs)

    # Slide 22: Best Practices - ECS
    create_ecs_best_practices_slide(prs)

    # Slide 23: Security & Compliance
    create_security_compliance_slide(prs)

    # Slide 24: Monitoring & Verification
    create_monitoring_verification_slide(prs)

    # Slide 25: Rollback Strategies
    create_rollback_strategies_slide(prs)

    # Slide 26: Implementation Timeline
    create_implementation_timeline_slide(prs)

    # Slide 27: Key Takeaways
    create_key_takeaways_slide(prs)

    # Slide 28: Next Steps
    create_next_steps_slide(prs)

    # Slide 29: Resources
    create_resources_slide(prs)

    # Slide 30: Thank You
    create_thank_you_slide(prs)

    return prs

def add_title_to_slide(slide, title_text):
    """Add a title to a slide"""
    # Add title as a text box since we're using blank layout
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_GRAY):
    """Add a text box to a slide"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return textbox

def add_bullet_points(slide, left, top, width, height, items, font_size=18):
    """Add bullet points to a slide"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(items):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)

def add_diagram_box(slide, left, top, width, height, text, color=PRIMARY_COLOR):
    """Add a diagram box to a slide"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color

    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return shape

def add_arrow(slide, x1, y1, x2, y2):
    """Add an arrow connector between two points"""
    connector = slide.shapes.add_connector(
        1,  # Straight connector
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = DARK_GRAY
    connector.line.width = Pt(2)
    return connector

# ============================================================================
# SLIDE CREATION FUNCTIONS
# ============================================================================

def create_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Main title
    title_box = add_textbox(slide, 1, 2.5, 8, 1,
                           "Harness Continuous Delivery (CD)",
                           font_size=54, bold=True, color=PRIMARY_COLOR)
    title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = add_textbox(slide, 1, 3.5, 8, 0.8,
                              "VM & ECS Deployment Workflows",
                              font_size=36, bold=False, color=SECONDARY_COLOR)
    subtitle_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Additional info
    info_box = add_textbox(slide, 1, 5, 8, 0.5,
                          "Detailed Architecture & Implementation Guide",
                          font_size=24, bold=False, color=DARK_GRAY)
    info_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_agenda_slide(prs):
    """Slide 2: Agenda"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "Agenda")

    items = [
        "1. Introduction to Harness CD",
        "2. CD Flow Overview",
        "3. VM Deployment Architecture & Workflow",
        "4. VM Deployment Pipeline Details",
        "5. ECS Deployment Architecture & Workflow",
        "6. ECS Deployment Strategies (Rolling, Blue-Green, Canary)",
        "7. Best Practices & Security",
        "8. Monitoring & Rollback Strategies",
        "9. Implementation Timeline & Next Steps"
    ]
    add_bullet_points(slide, 1, 1.5, 8, 5, items, font_size=20)

def create_intro_slide(prs):
    """Slide 3: Introduction to Harness CD"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "What is Harness CD?")

    # Definition
    add_textbox(slide, 1, 1.5, 8, 0.5,
               "Modern Continuous Delivery platform for automating application deployments",
               font_size=20, bold=True, color=PRIMARY_COLOR)

    # Key capabilities
    add_textbox(slide, 1, 2.2, 8, 0.4, "Key Capabilities:", font_size=18, bold=True)

    items = [
        "✓ Multi-cloud deployments (AWS, Azure, GCP, On-Premise)",
        "✓ Multiple deployment targets (VM, ECS, Kubernetes, Serverless)",
        "✓ Advanced deployment strategies (Rolling, Blue-Green, Canary)",
        "✓ AI-powered continuous verification",
        "✓ Automated rollback capabilities",
        "✓ Policy as Code for governance"
    ]
    add_bullet_points(slide, 1.5, 2.8, 7, 2.5, items, font_size=16)

    # Benefits
    add_textbox(slide, 1, 5.5, 8, 0.4, "Business Benefits:", font_size=18, bold=True)

    benefits = [
        "• 70% reduction in deployment time",
        "• 40-60% fewer production incidents",
        "• Zero-downtime deployments",
        "• Complete audit trails and compliance"
    ]
    add_bullet_points(slide, 1.5, 6, 7, 1.2, benefits, font_size=16)

def create_cd_flow_overview_slide(prs):
    """Slide 4: CD Flow Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "CD Flow Overview")

    # Flow diagram
    y_pos = 2
    x_start = 1.5
    box_width = 1.8
    box_height = 0.6
    spacing = 2.2

    # Row 1: Code to Artifact
    add_diagram_box(slide, x_start, y_pos, box_width, box_height, "Code\nCommit", PRIMARY_COLOR)
    add_arrow(slide, x_start + box_width, y_pos + 0.3, x_start + spacing, y_pos + 0.3)

    add_diagram_box(slide, x_start + spacing, y_pos, box_width, box_height, "CI\nPipeline", SECONDARY_COLOR)
    add_arrow(slide, x_start + spacing + box_width, y_pos + 0.3, x_start + spacing * 2, y_pos + 0.3)

    add_diagram_box(slide, x_start + spacing * 2, y_pos, box_width, box_height, "Artifact\nRegistry", ACCENT_COLOR)

    # Row 2: Harness CD
    y_pos = 3.5
    add_diagram_box(slide, 3, y_pos, 4, box_height, "Harness CD Pipeline", PRIMARY_COLOR)

    # Row 3: Environments
    y_pos = 4.8
    x_start = 0.5
    box_width = 1.5

    add_diagram_box(slide, x_start, y_pos, box_width, box_height, "Dev\n✓ Deploy", SECONDARY_COLOR)
    add_arrow(slide, x_start + box_width, y_pos + 0.3, x_start + box_width + 0.3, y_pos + 0.3)

    add_diagram_box(slide, x_start + 2, y_pos, box_width, box_height, "QA\n✓ Test", SECONDARY_COLOR)
    add_arrow(slide, x_start + 2 + box_width, y_pos + 0.3, x_start + 2 + box_width + 0.3, y_pos + 0.3)

    add_diagram_box(slide, x_start + 4, y_pos, box_width, box_height, "Staging\n✓ Verify", SECONDARY_COLOR)
    add_arrow(slide, x_start + 4 + box_width, y_pos + 0.3, x_start + 4 + box_width + 0.3, y_pos + 0.3)

    add_diagram_box(slide, x_start + 6, y_pos, 1.8, box_height, "Production\n✓ Monitor", ACCENT_COLOR)

def create_vm_architecture_slide(prs):
    """Slide 5: VM Deployment Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "VM Deployment Architecture")

    # Harness Platform
    add_diagram_box(slide, 2, 1.5, 6, 0.8, "Harness CD Platform", PRIMARY_COLOR)

    # Pipeline components
    y_pos = 2.5
    add_diagram_box(slide, 2.5, y_pos, 1.5, 0.5, "Pipeline", SECONDARY_COLOR)
    add_diagram_box(slide, 4.2, y_pos, 1.5, 0.5, "Service", SECONDARY_COLOR)
    add_diagram_box(slide, 5.9, y_pos, 1.5, 0.5, "Environment", SECONDARY_COLOR)

    # Delegate
    add_diagram_box(slide, 3.5, 3.3, 3, 0.6, "Harness Delegate (Agent)", ACCENT_COLOR)

    # Connection
    add_textbox(slide, 4.5, 4, 1, 0.3, "SSH/Agent", font_size=12, bold=True)
    add_arrow(slide, 5, 3.9, 5, 4.4)

    # VM Servers
    y_pos = 4.7
    add_diagram_box(slide, 1, y_pos, 2.3, 1, "VM Server 1\n(Production)\n• App Service\n• Port: 8080", SECONDARY_COLOR)
    add_diagram_box(slide, 3.8, y_pos, 2.3, 1, "VM Server 2\n(Staging)\n• App Service\n• Port: 8080", SECONDARY_COLOR)
    add_diagram_box(slide, 6.6, y_pos, 2.3, 1, "VM Server N\n(Dev/QA)\n• App Service\n• Port: 8080", SECONDARY_COLOR)

    # Legend
    add_textbox(slide, 0.5, 6.5, 9, 0.8,
               "🔧 SSH Connection  |  📦 Artifact Deployment  |  ✓ Health Checks  |  🔄 Rollback Support",
               font_size=14, color=DARK_GRAY)

def create_vm_dataflow_slide(prs):
    """Slide 6: VM Deployment Dataflow"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "VM Deployment Dataflow")

    flow_items = [
        "1. Trigger: Git Push → Manual → Scheduled",
        "2. Artifact Fetch: Docker Registry / Artifactory / S3",
        "3. Delegate Selection: Choose healthy delegate",
        "4. SSH Connection: Establish secure connection to VM",
        "5. Pre-Deployment:",
        "   • Backup current application",
        "   • Stop running services",
        "   • Validate environment (disk, memory, dependencies)",
        "6. Deployment Execution:",
        "   • Transfer artifacts via SCP/SFTP",
        "   • Extract and install application",
        "   • Update configuration files",
        "7. Post-Deployment:",
        "   • Start services",
        "   • Health checks (HTTP endpoints, process status)",
        "   • Smoke tests",
        "8. Continuous Verification: Monitor metrics & auto-rollback",
        "9. Completion: Success ✓ or Rollback 🔄"
    ]
    add_bullet_points(slide, 1, 1.5, 8, 5.5, flow_items, font_size=16)

def create_vm_components_slide(prs):
    """Slide 7: VM Deployment Components"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "VM Deployment Components")

    # Component boxes
    components = [
        ("Harness Delegate", "• Lightweight agent\n• Executes commands\n• Reports status", 1, 1.8),
        ("SSH Connection", "• Secure shell (Port 22)\n• SSH key authentication\n• Encrypted communication", 4, 1.8),
        ("VM Infrastructure", "• Linux/Windows OS\n• Network access required\n• Sufficient resources", 7, 1.8),
        ("Artifacts", "• TAR/ZIP/RPM/DEB\n• Docker images\n• Versioned & tagged", 1, 4),
        ("Configuration", "• App properties\n• Environment configs\n• Secrets management", 4, 4),
        ("Monitoring", "• Health checks\n• Prometheus/Datadog\n• Custom metrics", 7, 4)
    ]

    for title, content, x, y in components:
        add_diagram_box(slide, x, y, 2.3, 0.5, title, PRIMARY_COLOR)
        add_textbox(slide, x, y + 0.6, 2.3, 0.8, content, font_size=12)

def create_vm_pipeline_steps_slide(prs):
    """Slide 8: VM Deployment Pipeline Steps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "VM Pipeline: Detailed Steps")

    steps = [
        "Step 1: Backup Current Application",
        "   • Create timestamped backup directory",
        "   • Copy application files",
        "   • Archive for rollback",
        "",
        "Step 2: Stop Services",
        "   • systemctl stop myapp",
        "   • Verify service stopped",
        "",
        "Step 3: Download & Extract Artifacts",
        "   • Fetch from artifact repository",
        "   • Extract to deployment directory",
        "   • Set proper permissions",
        "",
        "Step 4: Update Configuration",
        "   • Copy environment-specific configs",
        "   • Replace variables",
        "   • Validate configuration",
        "",
        "Step 5: Start Services",
        "   • systemctl start myapp",
        "   • Verify service health"
    ]
    add_bullet_points(slide, 1, 1.5, 8, 5.5, steps, font_size=14)

def create_vm_pre_deployment_slide(prs):
    """Slide 9: VM Pre-Deployment Phase"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "VM Pre-Deployment Phase")

    # Phase diagram
    add_diagram_box(slide, 3, 1.5, 4, 0.6, "Pre-Deployment Phase", PRIMARY_COLOR)

    phases = [
        ("Environment\nValidation", 1.5, 2.5, "• Disk space\n• Memory\n• Dependencies"),
        ("Backup\nCurrent App", 4, 2.5, "• Create backup\n• Archive\n• Store safely"),
        ("Stop\nServices", 6.5, 2.5, "• Stop app\n• Verify stopped\n• Wait for drain")
    ]

    for title, x, y, details in phases:
        add_diagram_box(slide, x, y, 2, 0.6, title, SECONDARY_COLOR)
        add_textbox(slide, x, y + 0.7, 2, 0.8, details, font_size=12)

    # Best practices
    add_textbox(slide, 1, 4.5, 8, 0.4, "Best Practices:", font_size=18, bold=True)

    practices = [
        "✓ Always backup before deployment",
        "✓ Validate environment prerequisites",
        "✓ Check for running processes",
        "✓ Verify sufficient disk space (at least 2x artifact size)",
        "✓ Document pre-deployment state"
    ]
    add_bullet_points(slide, 1.5, 5, 7, 2, practices, font_size=14)

def create_vm_deployment_execution_slide(prs):
    """Slide 10: VM Deployment Execution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "VM Deployment Execution")

    # Execution flow
    add_diagram_box(slide, 3, 1.5, 4, 0.6, "Deployment Execution Phase", PRIMARY_COLOR)

    steps = [
        ("Transfer\nArtifacts", 1, 2.5, "SCP/SFTP\nVerify transfer"),
        ("Extract &\nInstall", 3, 2.5, "Unzip/Untar\nSet permissions"),
        ("Update\nConfig", 5, 2.5, "Copy configs\nReplace vars"),
        ("Start\nServices", 7, 2.5, "Start app\nVerify running")
    ]

    for i, (title, x, y, details) in enumerate(steps):
        add_diagram_box(slide, x, y, 1.7, 0.6, title, SECONDARY_COLOR)
        add_textbox(slide, x, y + 0.7, 1.7, 0.6, details, font_size=11)

        if i < len(steps) - 1:
            add_arrow(slide, x + 1.7, y + 0.3, x + 2, y + 0.3)

    # Example commands
    add_textbox(slide, 1, 4.2, 8, 0.4, "Example Commands:", font_size=16, bold=True)

    commands = [
        "# Transfer artifact",
        "scp myapp.tar.gz user@vm:/opt/myapp/",
        "",
        "# Extract and install",
        "tar -xzf myapp.tar.gz && chmod +x bin/myapp",
        "",
        "# Update configuration",
        "cp config/app.properties /opt/myapp/conf/",
        "",
        "# Start service",
        "systemctl start myapp && systemctl status myapp"
    ]
    add_bullet_points(slide, 1.5, 4.7, 7, 2.5, commands, font_size=12)

def create_vm_post_deployment_slide(prs):
    """Slide 11: VM Post-Deployment Phase"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "VM Post-Deployment Phase")

    add_diagram_box(slide, 3, 1.5, 4, 0.6, "Post-Deployment Phase", PRIMARY_COLOR)

    steps = [
        ("Health\nChecks", 1.5, 2.5, "• Process status\n• Port listening\n• HTTP /health"),
        ("Smoke\nTests", 4, 2.5, "• API endpoints\n• Functionality\n• Error logs"),
        ("Continuous\nVerification", 6.5, 2.5, "• Metrics\n• Performance\n• Auto-rollback")
    ]

    for title, x, y, details in steps:
        add_diagram_box(slide, x, y, 2, 0.6, title, SECONDARY_COLOR)
        add_textbox(slide, x, y + 0.7, 2, 0.8, details, font_size=12)

    # Verification checklist
    add_textbox(slide, 1, 4.5, 8, 0.4, "Verification Checklist:", font_size=16, bold=True)

    checklist = [
        "✓ Service is running (systemctl status)",
        "✓ Application port is listening (netstat -tulpn)",
        "✓ Health endpoint returns 200 OK",
        "✓ API endpoints respond correctly",
        "✓ No errors in application logs",
        "✓ Response times within acceptable range",
        "✓ Resource usage (CPU, memory) normal",
        "✓ External dependencies accessible"
    ]
    add_bullet_points(slide, 1.5, 5, 7, 2.2, checklist, font_size=13)

def create_ecs_architecture_slide(prs):
    """Slide 12: ECS Deployment Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "ECS Deployment Architecture")

    # Harness Platform
    add_diagram_box(slide, 2, 1.5, 6, 0.7, "Harness CD Platform", PRIMARY_COLOR)

    # Pipeline components
    y_pos = 2.4
    add_diagram_box(slide, 2.5, y_pos, 1.5, 0.5, "Pipeline", SECONDARY_COLOR)
    add_diagram_box(slide, 4.2, y_pos, 1.5, 0.5, "Service", SECONDARY_COLOR)
    add_diagram_box(slide, 5.9, y_pos, 1.5, 0.5, "Environment", SECONDARY_COLOR)

    # Delegate
    add_diagram_box(slide, 3.5, 3.2, 3, 0.5, "Harness Delegate (AWS VPC)", ACCENT_COLOR)

    # AWS API
    add_textbox(slide, 4.3, 3.9, 1.4, 0.3, "AWS APIs", font_size=12, bold=True)
    add_arrow(slide, 5, 3.7, 5, 4.3)

    # ECS Components
    y_pos = 4.6
    add_diagram_box(slide, 1, y_pos, 2, 0.8, "ECS Cluster\n(Production)", SECONDARY_COLOR)
    add_diagram_box(slide, 3.5, y_pos, 2, 0.8, "ECS Service\n(myapp-service)", SECONDARY_COLOR)
    add_diagram_box(slide, 6, y_pos, 2, 0.8, "Task Definition\n(myapp-task)", SECONDARY_COLOR)

    # ALB and ECR
    y_pos = 5.8
    add_diagram_box(slide, 2, y_pos, 2.5, 0.6, "Application Load\nBalancer (ALB)", ACCENT_COLOR)
    add_diagram_box(slide, 5.5, y_pos, 2.5, 0.6, "ECR Registry\n(Container Images)", ACCENT_COLOR)

    # Legend
    add_textbox(slide, 0.5, 6.8, 9, 0.6,
               "☁️ AWS Services  |  📦 Container Deployment  |  ⚖️ Load Balancing  |  🔄 Auto Scaling",
               font_size=13, color=DARK_GRAY)

def create_ecs_dataflow_slide(prs):
    """Slide 13: ECS Deployment Dataflow"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "ECS Deployment Dataflow")

    flow_items = [
        "1. Trigger: CI Pipeline Complete → Manual → Scheduled",
        "2. ECR Artifact Fetch:",
        "   • Authenticate with AWS ECR",
        "   • Pull image metadata (tag, digest, size)",
        "   • Validate image exists",
        "3. AWS Connection: IAM role/keys → Validate permissions",
        "4. Pre-Deployment:",
        "   • Get current task definition revision",
        "   • Get current service status",
        "   • Validate ECS cluster capacity",
        "5. Create New Task Definition:",
        "   • Update container image",
        "   • Update environment variables",
        "   • Register new revision",
        "6. Deployment Strategy: Rolling / Blue-Green / Canary",
        "7. ECS Service Update:",
        "   • Start new tasks with new task definition",
        "   • Register tasks with ALB target group",
        "   • Wait for health checks to pass",
        "   • Deregister and stop old tasks",
        "8. Verification: Task status, health checks, endpoints",
        "9. Completion: Success ✓ or Auto-Rollback 🔄"
    ]
    add_bullet_points(slide, 1, 1.5, 8, 5.5, flow_items, font_size=14)

def create_ecs_components_slide(prs):
    """Slide 14: ECS Components"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "Amazon ECS Components")

    components = [
        ("ECS Cluster", "• Logical grouping\n• EC2 or Fargate\n• Resource management", 1, 1.8),
        ("ECS Service", "• Maintains task count\n• Load balancing\n• Auto scaling", 4, 1.8),
        ("Task Definition", "• Container blueprint\n• CPU/Memory\n• Image & ports", 7, 1.8),
        ("Tasks", "• Running instances\n• Fargate/EC2\n• VPC networking", 1, 4),
        ("ALB/NLB", "• Traffic routing\n• Health checks\n• SSL termination", 4, 4),
        ("ECR", "• Image repository\n• Integrated with ECS\n• Image scanning", 7, 4)
    ]

    for title, content, x, y in components:
        add_diagram_box(slide, x, y, 2.3, 0.5, title, PRIMARY_COLOR)
        add_textbox(slide, x, y + 0.6, 2.3, 0.8, content, font_size=12)

def create_ecs_task_definition_slide(prs):
    """Slide 15: ECS Task Definition"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "ECS Task Definition")

    add_textbox(slide, 1, 1.5, 8, 0.4, "Task Definition Structure:", font_size=18, bold=True)

    # Task definition components
    components = [
        ("Container\nDefinitions", 1.5, 2.2, "• Image\n• Ports\n• Environment"),
        ("Resource\nLimits", 4, 2.2, "• CPU: 256\n• Memory: 512MB\n• GPU (optional)"),
        ("Networking", 6.5, 2.2, "• Network mode\n• awsvpc\n• Security groups")
    ]

    for title, x, y, details in components:
        add_diagram_box(slide, x, y, 2, 0.5, title, SECONDARY_COLOR)
        add_textbox(slide, x, y + 0.6, 2, 0.8, details, font_size=11)

    # Example configuration
    add_textbox(slide, 1, 3.8, 8, 0.4, "Example Task Definition:", font_size=16, bold=True)

    config = [
        "Family: myapp-task",
        "Launch Type: FARGATE",
        "CPU: 256 (0.25 vCPU)",
        "Memory: 512 MB",
        "",
        "Container Definition:",
        "  • Name: myapp-container",
        "  • Image: 123456789012.dkr.ecr.us-east-1.amazonaws.com/myapp:v1.2.3",
        "  • Port Mapping: 8080:8080",
        "  • Environment Variables: ENVIRONMENT=production",
        "  • Log Configuration: awslogs → CloudWatch",
        "",
        "Network Mode: awsvpc",
        "Requires Compatibilities: FARGATE"
    ]
    add_bullet_points(slide, 1.5, 4.3, 7, 2.8, config, font_size=12)

def create_ecs_deployment_strategies_slide(prs):
    """Slide 16: ECS Deployment Strategies"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "ECS Deployment Strategies")

    strategies = [
        ("Rolling Update", 1.5, 2, SECONDARY_COLOR,
         "• Default strategy\n• Updates one by one\n• Minimal downtime\n• Simple rollback",
         "Use When:\n• Standard deployments\n• Low-risk changes"),

        ("Blue-Green", 4, 2, PRIMARY_COLOR,
         "• Two environments\n• Instant switch\n• Zero downtime\n• Easy rollback",
         "Use When:\n• Critical apps\n• Need instant rollback"),

        ("Canary", 6.5, 2, ACCENT_COLOR,
         "• Gradual rollout\n• 10% → 50% → 100%\n• Risk mitigation\n• Performance validation",
         "Use When:\n• High-risk changes\n• Need validation")
    ]

    for title, x, y, color, features, use_case in strategies:
        add_diagram_box(slide, x, y, 2, 0.5, title, color)
        add_textbox(slide, x, y + 0.6, 2, 1.5, features, font_size=11)
        add_textbox(slide, x, y + 2.2, 2, 0.8, use_case, font_size=10, bold=True, color=color)

def create_ecs_rolling_deployment_slide(prs):
    """Slide 17: ECS Rolling Deployment"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "ECS Rolling Deployment")

    add_textbox(slide, 1, 1.5, 8, 0.4, "Rolling Update Process:", font_size=18, bold=True)

    # Process flow
    steps = [
        ("Current State", 1.2, 2.2, "Old tasks: 3\nRunning v1.0"),
        ("Start New Task", 3.2, 2.2, "New task: 1\nStarting v1.1"),
        ("Health Check", 5.2, 2.2, "New task: 1\nHealthy ✓"),
        ("Stop Old Task", 7.2, 2.2, "Old task: 2\nNew task: 1")
    ]

    for i, (title, x, y, details) in enumerate(steps):
        add_diagram_box(slide, x, y, 1.6, 0.5, title, SECONDARY_COLOR)
        add_textbox(slide, x, y + 0.6, 1.6, 0.6, details, font_size=10)

        if i < len(steps) - 1:
            add_arrow(slide, x + 1.6, y + 0.25, x + 1.9, y + 0.25)

    # Configuration
    add_textbox(slide, 1, 3.8, 8, 0.4, "Deployment Configuration:", font_size=16, bold=True)

    config = [
        "Deployment Configuration:",
        "  • Maximum Percent: 200%",
        "    (Can run up to 6 tasks during deployment if desired count is 3)",
        "  • Minimum Healthy Percent: 100%",
        "    (Must maintain at least 3 healthy tasks)",
        "",
        "Deployment Steps:",
        "  1. Start new task with new task definition",
        "  2. Wait for task to reach RUNNING state",
        "  3. Register task with target group",
        "  4. Wait for health checks to pass",
        "  5. Deregister old task from target group",
        "  6. Stop old task",
        "  7. Repeat for remaining tasks"
    ]
    add_bullet_points(slide, 1.5, 4.3, 7, 2.5, config, font_size=13)

def create_ecs_blue_green_slide(prs):
    """Slide 18: ECS Blue-Green Deployment"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "ECS Blue-Green Deployment")

    add_textbox(slide, 1, 1.5, 8, 0.4, "Blue-Green Deployment Process:", font_size=18, bold=True)

    # Process visualization
    phases = [
        ("Blue (Old)", 1.5, 2.2, PRIMARY_COLOR, "100% traffic\nv1.0 running"),
        ("Green (New)", 5, 2.2, SECONDARY_COLOR, "0% traffic\nv1.1 deploying"),
        ("Switch", 1.5, 3.5, ACCENT_COLOR, "0% traffic\nv1.0 standby"),
        ("Complete", 5, 3.5, SECONDARY_COLOR, "100% traffic\nv1.1 running")
    ]

    for title, x, y, color, details in phases:
        add_diagram_box(slide, x, y, 2.8, 0.5, title, color)
        add_textbox(slide, x, y + 0.6, 2.8, 0.6, details, font_size=12)

    # Traffic shifting
    add_textbox(slide, 1, 4.7, 8, 0.4, "Traffic Shifting Strategy:", font_size=16, bold=True)

    strategy = [
        "Phase 1: Deploy Green Environment",
        "  • Create new ECS service with new task definition",
        "  • Deploy to separate target group",
        "  • Run health checks and smoke tests",
        "",
        "Phase 2: Gradual Traffic Shift",
        "  • 10% traffic to Green → Verify 5 min",
        "  • 50% traffic to Green → Verify 5 min",
        "  • 100% traffic to Green → Monitor",
        "",
        "Phase 3: Cleanup",
        "  • Keep Blue environment for 24h (rollback window)",
        "  • Terminate Blue environment after validation"
    ]
    add_bullet_points(slide, 1.5, 5.2, 7, 2, strategy, font_size=12)

def create_ecs_canary_deployment_slide(prs):
    """Slide 19: ECS Canary Deployment"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "ECS Canary Deployment")

    add_textbox(slide, 1, 1.5, 8, 0.4, "Canary Deployment Process:", font_size=18, bold=True)

    # Canary phases
    phases = [
        ("10% Canary", 1.5, 2.2, ACCENT_COLOR, "1 new task\n9 old tasks\nVerify 5 min"),
        ("50% Canary", 4, 2.2, SECONDARY_COLOR, "5 new tasks\n5 old tasks\nVerify 10 min"),
        ("100% Deploy", 6.5, 2.2, PRIMARY_COLOR, "10 new tasks\n0 old tasks\nMonitor")
    ]

    for title, x, y, color, details in phases:
        add_diagram_box(slide, x, y, 2, 0.5, title, color)
        add_textbox(slide, x, y + 0.6, 2, 0.8, details, font_size=11)

    # Canary configuration
    add_textbox(slide, 1, 3.8, 8, 0.4, "Canary Configuration:", font_size=16, bold=True)

    config = [
        "Step 1: Deploy Canary (10%)",
        "  • Instance Selection: Count = 1",
        "  • Deploy new task definition to 1 task",
        "  • Verification Duration: 5 minutes",
        "  • Metrics: Error rate < 1%, Response time < baseline + 10%",
        "",
        "Step 2: Increase to 50%",
        "  • Instance Selection: Count = 5",
        "  • Verification Duration: 10 minutes",
        "  • Advanced metrics: CPU, memory, custom app metrics",
        "",
        "Step 3: Full Deployment",
        "  • Delete canary tasks",
        "  • Rolling deployment to remaining tasks",
        "  • Continuous monitoring for 30 minutes",
        "",
        "Auto-Rollback Triggers:",
        "  • Error rate exceeds threshold",
        "  • Health checks fail",
        "  • Performance degradation detected"
    ]
    add_bullet_points(slide, 1.5, 4.3, 7, 2.8, config, font_size=11)

def create_deployment_comparison_slide(prs):
    """Slide 20: Deployment Comparison"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "Deployment Strategy Comparison")

    # Table headers
    add_textbox(slide, 1, 1.5, 2, 0.4, "Strategy", font_size=14, bold=True, color=PRIMARY_COLOR)
    add_textbox(slide, 3, 1.5, 1.5, 0.4, "Downtime", font_size=14, bold=True, color=PRIMARY_COLOR)
    add_textbox(slide, 4.5, 1.5, 1.5, 0.4, "Rollback", font_size=14, bold=True, color=PRIMARY_COLOR)
    add_textbox(slide, 6, 1.5, 1.5, 0.4, "Complexity", font_size=14, bold=True, color=PRIMARY_COLOR)
    add_textbox(slide, 7.5, 1.5, 1.5, 0.4, "Best For", font_size=14, bold=True, color=PRIMARY_COLOR)

    # Rolling Update
    y = 2.1
    add_textbox(slide, 1, y, 2, 0.3, "Rolling Update", font_size=12, bold=True)
    add_textbox(slide, 3, y, 1.5, 0.3, "Minimal", font_size=12)
    add_textbox(slide, 4.5, y, 1.5, 0.3, "Manual", font_size=12)
    add_textbox(slide, 6, y, 1.5, 0.3, "Low", font_size=12)
    add_textbox(slide, 7.5, y, 1.5, 0.3, "Standard", font_size=11)

    # Blue-Green
    y = 2.6
    add_textbox(slide, 1, y, 2, 0.3, "Blue-Green", font_size=12, bold=True)
    add_textbox(slide, 3, y, 1.5, 0.3, "Zero", font_size=12)
    add_textbox(slide, 4.5, y, 1.5, 0.3, "Instant", font_size=12)
    add_textbox(slide, 6, y, 1.5, 0.3, "Medium", font_size=12)
    add_textbox(slide, 7.5, y, 1.5, 0.3, "Critical apps", font_size=11)

    # Canary
    y = 3.1
    add_textbox(slide, 1, y, 2, 0.3, "Canary", font_size=12, bold=True)
    add_textbox(slide, 3, y, 1.5, 0.3, "Zero", font_size=12)
    add_textbox(slide, 4.5, y, 1.5, 0.3, "Automatic", font_size=12)
    add_textbox(slide, 6, y, 1.5, 0.3, "High", font_size=12)
    add_textbox(slide, 7.5, y, 1.5, 0.3, "High-risk", font_size=11)

    # Decision matrix
    add_textbox(slide, 1, 4, 8, 0.4, "Decision Matrix:", font_size=16, bold=True)

    matrix = [
        "Choose Rolling Update when:",
        "  • Standard, low-risk deployments",
        "  • Simple application architecture",
        "  • Limited infrastructure capacity",
        "",
        "Choose Blue-Green when:",
        "  • Zero downtime is critical",
        "  • Need instant rollback capability",
        "  • Sufficient infrastructure capacity (2x)",
        "",
        "Choose Canary when:",
        "  • High-risk changes (major updates, architecture changes)",
        "  • Need gradual validation with real traffic",
        "  • Want automated rollback based on metrics"
    ]
    add_bullet_points(slide, 1.5, 4.5, 7, 2.5, matrix, font_size=12)

def create_vm_best_practices_slide(prs):
    """Slide 21: Best Practices - VM"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "VM Deployment Best Practices")

    practices = [
        "SSH Security:",
        "  ✓ Use SSH keys instead of passwords",
        "  ✓ Rotate SSH keys regularly (every 90 days)",
        "  ✓ Use bastion hosts for production access",
        "  ✓ Implement SSH key management solution",
        "",
        "Backup Strategy:",
        "  ✓ Always backup before deployment",
        "  ✓ Keep multiple backup versions (last 5 deployments)",
        "  ✓ Test restore procedures quarterly",
        "  ✓ Store backups off-server (S3, NAS)",
        "",
        "Health Checks:",
        "  ✓ Implement comprehensive health checks",
        "  ✓ Check both process status and HTTP endpoints",
        "  ✓ Set appropriate timeouts (30-60 seconds)",
        "  ✓ Retry failed health checks (3 attempts)",
        "",
        "Rollback Plan:",
        "  ✓ Test rollback procedures in staging",
        "  ✓ Keep previous versions available",
        "  ✓ Document rollback steps in runbook",
        "  ✓ Automate rollback when possible"
    ]
    add_bullet_points(slide, 1, 1.5, 8, 5.5, practices, font_size=14)

def create_ecs_best_practices_slide(prs):
    """Slide 22: Best Practices - ECS"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "ECS Deployment Best Practices")

    practices = [
        "Task Definition Management:",
        "  ✓ Version control task definitions in Git",
        "  ✓ Use parameterized task definitions",
        "  ✓ Test in non-production first",
        "  ✓ Keep previous revisions (minimum 10)",
        "",
        "Resource Management:",
        "  ✓ Right-size CPU and memory (monitor and adjust)",
        "  ✓ Use Fargate Spot for cost savings (non-prod)",
        "  ✓ Implement auto-scaling based on metrics",
        "  ✓ Set resource reservations and limits",
        "",
        "Load Balancer Configuration:",
        "  ✓ Use Application Load Balancer (ALB) for HTTP/HTTPS",
        "  ✓ Configure health checks: /health endpoint",
        "  ✓ Health check interval: 30s, Timeout: 5s",
        "  ✓ Set up SSL/TLS termination at ALB",
        "",
        "Deployment Configuration:",
        "  ✓ Enable deployment circuit breaker with auto-rollback",
        "  ✓ Set health check grace period: 60-120 seconds",
        "  ✓ Use appropriate deployment strategy",
        "  ✓ Configure CloudWatch alarms for automated alerts"
    ]
    add_bullet_points(slide, 1, 1.5, 8, 5.5, practices, font_size=13)

def create_security_compliance_slide(prs):
    """Slide 23: Security & Compliance"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "Security & Compliance")

    # Security features
    add_diagram_box(slide, 1.5, 1.5, 3, 0.6, "🔒 Secrets Management", SECONDARY_COLOR)
    add_diagram_box(slide, 5.5, 1.5, 3, 0.6, "👥 Access Control (RBAC)", SECONDARY_COLOR)

    security_items = [
        "Secrets Management:",
        "  • Native secrets management in Harness",
        "  • Integration with AWS Secrets Manager, HashiCorp Vault",
        "  • Encrypted at rest and in transit",
        "  • Automatic secret rotation support",
        "",
        "Access Control (RBAC):",
        "  • Role-based access control",
        "  • Fine-grained permissions",
        "  • Audit logs for all actions",
        "  • SSO integration (SAML, OAuth)",
        "",
        "Policy as Code (OPA):",
        "  • Enforce deployment policies",
        "  • Compliance checks before deployment",
        "  • Custom policies (e.g., require approval for prod)",
        "  • Policy violations prevent deployment",
        "",
        "Compliance:",
        "  • Complete audit trails",
        "  • SOC 2, ISO 27001 compliant",
        "  • GDPR compliance support",
        "  • Automated compliance reports"
    ]
    add_bullet_points(slide, 1, 2.3, 8, 4.5, security_items, font_size=13)

def create_monitoring_verification_slide(prs):
    """Slide 24: Monitoring & Verification"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "Monitoring & Continuous Verification")

    # Monitoring tools
    tools = [
        ("Prometheus", 1.5, 1.8, "• Metrics collection\n• Time-series DB\n• PromQL queries"),
        ("Datadog", 4.5, 1.8, "• APM monitoring\n• Infrastructure\n• Custom metrics"),
        ("CloudWatch", 7, 1.8, "• AWS native\n• Logs & metrics\n• Alarms")
    ]

    for title, x, y, details in tools:
        add_diagram_box(slide, x, y, 2, 0.5, title, SECONDARY_COLOR)
        add_textbox(slide, x, y + 0.6, 2, 0.7, details, font_size=11)

    # Verification metrics
    add_textbox(slide, 1, 3, 8, 0.4, "Continuous Verification Metrics:", font_size=16, bold=True)

    metrics = [
        "Error Rates:",
        "  • HTTP 5xx errors < 1%",
        "  • Application exceptions < 0.5%",
        "  • Failed requests < 0.1%",
        "",
        "Performance Metrics:",
        "  • Response time (P95) < baseline + 10%",
        "  • Response time (P99) < baseline + 20%",
        "  • Throughput ≥ baseline - 5%",
        "",
        "Resource Utilization:",
        "  • CPU utilization < 80%",
        "  • Memory utilization < 85%",
        "  • Disk I/O within normal range",
        "",
        "Auto-Rollback Triggers:",
        "  • Error rate exceeds threshold for 5 minutes",
        "  • Response time degradation > 50%",
        "  • Health check failures"
    ]
    add_bullet_points(slide, 1.5, 3.5, 7, 3.5, metrics, font_size=12)

def create_rollback_strategies_slide(prs):
    """Slide 25: Rollback Strategies"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "Rollback Strategies")

    # Rollback types
    add_diagram_box(slide, 1.5, 1.5, 3, 0.6, "Manual Rollback", PRIMARY_COLOR)
    add_diagram_box(slide, 5.5, 1.5, 3, 0.6, "Automatic Rollback", ACCENT_COLOR)

    strategies = [
        "Manual Rollback:",
        "  • Triggered by operator",
        "  • Review deployment status",
        "  • Execute rollback pipeline",
        "  • Verify rollback success",
        "",
        "VM Rollback Process:",
        "  1. Stop current application",
        "  2. Restore from backup (latest good version)",
        "  3. Start restored application",
        "  4. Verify health checks",
        "",
        "ECS Rollback Process:",
        "  1. Update ECS service with previous task definition",
        "  2. ECS automatically starts old tasks",
        "  3. ECS stops new tasks",
        "  4. Verify service health",
        "",
        "Automatic Rollback (AI-Powered):",
        "  • Continuous verification monitors metrics",
        "  • ML baseline comparison",
        "  • Anomaly detection",
        "  • Auto-rollback if degradation detected",
        "  • Notification sent to team",
        "",
        "Rollback Best Practices:",
        "  ✓ Test rollback procedures regularly",
        "  ✓ Keep rollback window ≤ 5 minutes",
        "  ✓ Document rollback procedures",
        "  ✓ Monitor post-rollback metrics"
    ]
    add_bullet_points(slide, 1, 2.3, 8, 4.8, strategies, font_size=12)

def create_implementation_timeline_slide(prs):
    """Slide 26: Implementation Timeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "Implementation Timeline")

    timeline = [
        "Phase 1: Planning & Setup (Week 1-2)",
        "  • Requirements gathering",
        "  • Harness account setup",
        "  • Connector configuration (Git, AWS, Docker)",
        "  • Delegate installation",
        "",
        "Phase 2: Environment Setup (Week 2-3)",
        "  • Create environments (Dev, QA, Staging, Prod)",
        "  • Define infrastructure (VM, ECS clusters)",
        "  • Configure secrets management",
        "  • Set up RBAC",
        "",
        "Phase 3: Pipeline Development (Week 3-5)",
        "  • Create service definitions",
        "  • Build deployment pipelines",
        "  • Configure deployment strategies",
        "  • Add verification steps",
        "",
        "Phase 4: Testing & Validation (Week 5-6)",
        "  • Test in Dev/QA environments",
        "  • Validate rollback procedures",
        "  • Performance testing",
        "  • Security validation",
        "",
        "Phase 5: Production Deployment (Week 7-8)",
        "  • Production readiness review",
        "  • Production deployment",
        "  • Monitoring and optimization",
        "  • Team training & documentation",
        "",
        "Total Timeline: 6-8 weeks for complete implementation"
    ]
    add_bullet_points(slide, 1, 1.5, 8, 5.5, timeline, font_size=13)

def create_key_takeaways_slide(prs):
    """Slide 27: Key Takeaways"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "Key Takeaways")

    takeaways = [
        "VM Deployments:",
        "  ✓ SSH-based deployments to traditional servers",
        "  ✓ Backup → Stop → Deploy → Start → Verify",
        "  ✓ Suitable for legacy applications and traditional infrastructure",
        "  ✓ Simple rollback via backup restoration",
        "",
        "ECS Deployments:",
        "  ✓ Container-based deployments on AWS",
        "  ✓ Multiple strategies: Rolling, Blue-Green, Canary",
        "  ✓ Zero-downtime deployments possible",
        "  ✓ Native AWS integration and auto-scaling",
        "",
        "Deployment Strategies:",
        "  • Rolling: Simple, minimal downtime",
        "  • Blue-Green: Zero downtime, instant rollback",
        "  • Canary: Gradual rollout, risk mitigation",
        "",
        "Critical Success Factors:",
        "  ✓ Comprehensive testing in non-production",
        "  ✓ Automated health checks and verification",
        "  ✓ Tested rollback procedures",
        "  ✓ Continuous monitoring and alerting",
        "  ✓ Team training and documentation",
        "",
        "Benefits:",
        "  • 70% faster deployments",
        "  • 40-60% fewer production incidents",
        "  • Complete visibility and audit trails",
        "  • Automated compliance and governance"
    ]
    add_bullet_points(slide, 1, 1.5, 8, 5.5, takeaways, font_size=13)

def create_next_steps_slide(prs):
    """Slide 28: Next Steps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "Next Steps")

    steps = [
        "Immediate Actions (This Week):",
        "  1. ✓ Set up Harness account (Free trial available)",
        "  2. ✓ Identify pilot application/service",
        "  3. ✓ Form project team (DevOps, Developer, QA)",
        "  4. ✓ Review existing deployment processes",
        "",
        "Short-term (Next 2 Weeks):",
        "  1. ✓ Install Harness delegate",
        "  2. ✓ Configure connectors (Git, Cloud, Artifact repos)",
        "  3. ✓ Create Dev environment",
        "  4. ✓ Build first deployment pipeline",
        "  5. ✓ Test deployment in Dev",
        "",
        "Medium-term (Next 4-6 Weeks):",
        "  1. ✓ Expand to QA and Staging environments",
        "  2. ✓ Implement deployment strategies",
        "  3. ✓ Set up continuous verification",
        "  4. ✓ Configure monitoring and alerting",
        "  5. ✓ Test rollback procedures",
        "",
        "Long-term (Next 2-3 Months):",
        "  1. ✓ Production deployment",
        "  2. ✓ Onboard additional applications",
        "  3. ✓ Team training and knowledge transfer",
        "  4. ✓ Continuous optimization",
        "",
        "Resources Needed:",
        "  • 1-2 DevOps Engineers",
        "  • 0.5-1 Developer",
        "  • 0.25-0.5 QA Engineer",
        "  • Access to infrastructure (VMs, AWS accounts)"
    ]
    add_bullet_points(slide, 1, 1.5, 8, 5.5, steps, font_size=12)

def create_resources_slide(prs):
    """Slide 29: Resources"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_to_slide(slide, "Resources & Documentation")

    resources = [
        "Official Documentation:",
        "  • Harness CD Docs: https://developer.harness.io/docs/continuous-delivery",
        "  • VM Deployments: https://developer.harness.io/docs/continuous-delivery/deploy-srv-diff-platforms/ssh-how-tos",
        "  • ECS Deployments: https://developer.harness.io/docs/continuous-delivery/deploy-srv-diff-platforms/aws-ecs-deployments",
        "",
        "Learning Resources:",
        "  • Harness University: https://university.harness.io",
        "  • Video Tutorials: YouTube - Harness Channel",
        "  • Community Forums: https://community.harness.io",
        "  • GitHub Examples: https://github.com/harness-community",
        "",
        "Internal Resources:",
        "  • Pipeline Templates: [Link to your internal templates]",
        "  • Deployment Runbooks: [Link to runbooks]",
        "  • Architecture Diagrams: [Link to architecture docs]",
        "  • Team Contact: [Your team Slack channel / Email]",
        "",
        "Support:",
        "  • Harness Support Portal: https://support.harness.io",
        "  • Professional Services: Available for complex implementations",
        "  • Training: Custom training sessions available",
        "",
        "Additional Files:",
        "  • HARNESS_CD_VM_ECS_WORKFLOW.md - Detailed technical guide",
        "  • CD_FLOW_TASK_BREAKDOWN.md - Complete task breakdown",
        "  • Configuration examples and YAML templates"
    ]
    add_bullet_points(slide, 1, 1.5, 8, 5.5, resources, font_size=11)

def create_thank_you_slide(prs):
    """Slide 30: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Thank you message
    thank_you_box = add_textbox(slide, 1, 2.5, 8, 1,
                                "Thank You!",
                                font_size=54, bold=True, color=PRIMARY_COLOR)
    thank_you_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Summary
    summary_box = add_textbox(slide, 1, 3.5, 8, 0.6,
                             "Questions?",
                             font_size=36, bold=False, color=SECONDARY_COLOR)
    summary_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Contact info
    contact_items = [
        "📧 Email: [your-email@company.com]",
        "💬 Slack: #harness-cd-support",
        "📚 Documentation: [Internal Wiki Link]",
        "🔗 Harness Platform: [Your Harness URL]"
    ]

    y_pos = 4.8
    for item in contact_items:
        contact_box = add_textbox(slide, 2, y_pos, 6, 0.3, item, font_size=18, color=DARK_GRAY)
        contact_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        y_pos += 0.4

# ============================================================================
# MAIN FUNCTION
# ============================================================================

def main():
    """Main function to generate the presentation"""
    print("Creating Harness CD Workflow Presentation...")

    try:
        prs = create_presentation()

        output_file = "/Users/tejasodanapalli/Desktop/Dailywork/Harness_CD_VM_ECS_Deployment_Workflows.pptx"
        prs.save(output_file)

        print(f"✓ Presentation created successfully!")
        print(f"✓ File saved: {output_file}")
        print(f"✓ Total slides: 30")
        print("\nPresentation includes:")
        print("  • VM Deployment Architecture & Workflow")
        print("  • ECS Deployment Architecture & Workflow")
        print("  • Deployment Strategies (Rolling, Blue-Green, Canary)")
        print("  • Best Practices & Security")
        print("  • Monitoring & Rollback Strategies")
        print("  • Implementation Timeline")

    except Exception as e:
        print(f"✗ Error creating presentation: {str(e)}")
        raise

if __name__ == "__main__":
    main()
