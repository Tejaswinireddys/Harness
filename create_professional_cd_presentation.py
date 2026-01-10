#!/usr/bin/env python3
"""
Professional Harness CD Architecture Presentation
Customer-sharable, comprehensive, and visually clear
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

# Professional Color Palette
HARNESS_BLUE = RGBColor(0, 125, 240)        # Primary brand color
DARK_BLUE = RGBColor(0, 51, 160)            # Secondary
SUCCESS_GREEN = RGBColor(34, 197, 94)       # Success indicators
WARNING_ORANGE = RGBColor(251, 146, 60)     # Warnings/Important
ERROR_RED = RGBColor(239, 68, 68)           # Errors
LIGHT_GRAY = RGBColor(243, 244, 246)        # Background
MEDIUM_GRAY = RGBColor(156, 163, 175)       # Secondary text
DARK_GRAY = RGBColor(55, 65, 81)            # Primary text
WHITE = RGBColor(255, 255, 255)

def create_presentation():
    """Create the comprehensive presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Create all slides
    create_title_slide(prs)
    create_agenda_slide(prs)
    create_introduction_slide(prs)
    create_architecture_overview_slide(prs)
    create_deployment_strategies_overview_slide(prs)

    # VM Deployment Section
    create_vm_architecture_overview_slide(prs)
    create_vm_detailed_architecture_slide(prs)
    create_vm_deployment_flow_phase1_slide(prs)
    create_vm_deployment_flow_phase2_slide(prs)
    create_vm_deployment_flow_phase3_slide(prs)
    create_vm_deployment_pipeline_slide(prs)
    create_vm_rollback_slide(prs)

    # ECS Deployment Section
    create_ecs_architecture_overview_slide(prs)
    create_ecs_detailed_architecture_slide(prs)
    create_ecs_deployment_flow_phase1_slide(prs)
    create_ecs_deployment_flow_phase2_slide(prs)
    create_ecs_blue_green_detailed_slide(prs)
    create_ecs_canary_detailed_slide(prs)
    create_ecs_pipeline_slide(prs)

    # Network & Security
    create_network_architecture_slide(prs)
    create_security_architecture_slide(prs)
    create_monitoring_slide(prs)

    # Best Practices & Implementation
    create_deployment_comparison_slide(prs)
    create_best_practices_slide(prs)
    create_implementation_roadmap_slide(prs)
    create_benefits_slide(prs)

    # Closing
    create_key_takeaways_slide(prs)
    create_next_steps_slide(prs)
    create_thank_you_slide(prs)

    return prs

def add_slide_with_title(prs, title_text):
    """Add a blank slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = HARNESS_BLUE

    return slide

def add_text_box(slide, left, top, width, height, text, font_size=16, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    """Add a text box to slide"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box

def add_shape_box(slide, left, top, width, height, text, bg_color, text_color=WHITE, font_size=14, bold=True):
    """Add a colored shape with text"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = bg_color

    frame = shape.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    return shape

def add_arrow(slide, x1, y1, x2, y2, color=MEDIUM_GRAY):
    """Add an arrow connector"""
    connector = slide.shapes.add_connector(
        1,  # Straight connector
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2.5)
    return connector

def add_bullet_list(slide, left, top, width, height, items, font_size=16, indent_level=0):
    """Add a bullet list"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = frame.paragraphs[0]
        else:
            p = frame.add_paragraph()

        p.text = item
        p.level = indent_level if isinstance(indent_level, int) else indent_level.get(i, 0)
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)

    return box

# =============================================================================
# SLIDE CREATION FUNCTIONS
# =============================================================================

def create_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = HARNESS_BLUE

    # Main title
    title_box = add_text_box(slide, 1, 2.5, 11.3, 1.5,
                             "Harness Continuous Delivery",
                             font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    subtitle_box = add_text_box(slide, 1, 4, 11.3, 0.8,
                                "End-to-End Architecture for VM & ECS Deployments",
                                font_size=32, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    # Footer
    footer_box = add_text_box(slide, 1, 6.5, 11.3, 0.5,
                              "Comprehensive Architecture Guide | 2026",
                              font_size=18, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

def create_agenda_slide(prs):
    """Slide 2: Agenda"""
    slide = add_slide_with_title(prs, "Agenda")

    items = [
        "Introduction to Harness CD",
        "Architecture Overview & Components",
        "Deployment Strategies",
        "",
        "VM Deployment Architecture",
        "    • Detailed architecture and components",
        "    • End-to-end deployment flow",
        "    • Pipeline configuration",
        "",
        "ECS Deployment Architecture",
        "    • Container orchestration architecture",
        "    • Blue-Green & Canary strategies",
        "    • AWS integration details",
        "",
        "Network & Security Architecture",
        "Monitoring & Continuous Verification",
        "Best Practices & Implementation Roadmap"
    ]

    add_bullet_list(slide, 1.5, 1.3, 10, 5.8, items, font_size=18)

def create_introduction_slide(prs):
    """Slide 3: Introduction"""
    slide = add_slide_with_title(prs, "What is Harness CD?")

    # Definition box
    add_shape_box(slide, 1.5, 1.5, 10, 0.6,
                  "Modern Continuous Delivery platform that automates software deployments across any environment",
                  HARNESS_BLUE, WHITE, font_size=18, bold=False)

    # Key capabilities
    add_text_box(slide, 1.5, 2.4, 10, 0.4, "Key Capabilities:", font_size=22, bold=True, color=HARNESS_BLUE)

    capabilities = [
        "✓ Multi-Cloud & Hybrid Deployments (AWS, Azure, GCP, On-Premise)",
        "✓ Multiple Deployment Targets (VMs, Containers, Kubernetes, Serverless)",
        "✓ Intelligent Deployment Strategies (Rolling, Blue-Green, Canary)",
        "✓ AI-Powered Continuous Verification with Auto-Rollback",
        "✓ Enterprise-Grade Security (RBAC, Secrets Management, Audit Trails)",
        "✓ GitOps & Infrastructure as Code Support"
    ]
    add_bullet_list(slide, 1.8, 2.9, 9.5, 2, capabilities, font_size=16)

    # Benefits section
    add_text_box(slide, 1.5, 5.2, 10, 0.4, "Business Impact:", font_size=22, bold=True, color=SUCCESS_GREEN)

    # Benefit boxes
    benefits = [
        ("70%", "Faster\nDeployments"),
        ("60%", "Fewer\nIncidents"),
        ("Zero", "Downtime\nReleases"),
        ("100%", "Audit\nCompliance")
    ]

    x_pos = 2
    for number, label in benefits:
        add_shape_box(slide, x_pos, 5.7, 2, 1.2, f"{number}\n{label}",
                     SUCCESS_GREEN, WHITE, font_size=18, bold=True)
        x_pos += 2.3

def create_architecture_overview_slide(prs):
    """Slide 4: Architecture Overview"""
    slide = add_slide_with_title(prs, "Harness CD: High-Level Architecture")

    y_start = 1.5

    # Control Plane
    add_shape_box(slide, 3, y_start, 7, 0.8, "Harness CD Control Plane (SaaS)",
                 HARNESS_BLUE, WHITE, font_size=20)

    # Components in Control Plane
    y_start += 1
    components = [
        "Pipeline\nOrchestrator",
        "Service\nManagement",
        "Environment\nManagement",
        "Secrets\nVault"
    ]
    x_pos = 2.5
    for comp in components:
        add_shape_box(slide, x_pos, y_start, 1.8, 0.7, comp,
                     DARK_BLUE, WHITE, font_size=11, bold=True)
        x_pos += 2

    # Arrow down
    add_arrow(slide, 6.5, y_start + 0.7, 6.5, y_start + 1.2, HARNESS_BLUE)
    add_text_box(slide, 5.5, y_start + 0.85, 2, 0.3,
                "WebSocket/HTTPS (TLS 1.3)", font_size=11, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    # Delegates
    y_start += 1.5
    add_shape_box(slide, 3.5, y_start, 6, 0.8, "Harness Delegates (Customer Network/VPC)",
                 WARNING_ORANGE, WHITE, font_size=18)

    # Delegate types
    y_start += 1
    delegates = ["Delegate 1\n(Active)", "Delegate 2\n(Active)", "Delegate 3\n(Standby)"]
    x_pos = 3.5
    for deleg in delegates:
        add_shape_box(slide, x_pos, y_start, 2, 0.6, deleg,
                     WARNING_ORANGE, WHITE, font_size=12)
        x_pos += 2.3

    # Arrows to targets
    add_arrow(slide, 4.5, y_start + 0.6, 3.5, y_start + 1.3, WARNING_ORANGE)
    add_arrow(slide, 6.5, y_start + 0.6, 6.5, y_start + 1.3, WARNING_ORANGE)
    add_arrow(slide, 8.5, y_start + 0.6, 9.5, y_start + 1.3, WARNING_ORANGE)

    # Target Infrastructure
    y_start += 1.5
    targets = [
        ("VM\nServers", DARK_BLUE),
        ("AWS ECS\nClusters", SUCCESS_GREEN),
        ("Kubernetes\nClusters", HARNESS_BLUE)
    ]
    x_pos = 2.5
    for target, color in targets:
        add_shape_box(slide, x_pos, y_start, 2.3, 0.7, target, color, WHITE, font_size=14)
        x_pos += 2.8

def create_deployment_strategies_overview_slide(prs):
    """Slide 5: Deployment Strategies Overview"""
    slide = add_slide_with_title(prs, "Deployment Strategies")

    strategies = [
        {
            "name": "Rolling Update",
            "color": HARNESS_BLUE,
            "icon": "🔄",
            "description": "Updates instances one at a time",
            "features": ["Minimal downtime", "Simple to implement", "Resource efficient"],
            "use_case": "Standard deployments, low-risk changes"
        },
        {
            "name": "Blue-Green",
            "color": SUCCESS_GREEN,
            "icon": "🔵🟢",
            "description": "Two identical environments for instant switch",
            "features": ["Zero downtime", "Instant rollback", "Full validation"],
            "use_case": "Critical apps, zero-downtime required"
        },
        {
            "name": "Canary",
            "color": WARNING_ORANGE,
            "icon": "🐤",
            "description": "Gradual rollout with traffic percentage",
            "features": ["Risk mitigation", "Incremental validation", "Auto-rollback"],
            "use_case": "High-risk changes, new features"
        }
    ]

    x_pos = 1.2
    for strategy in strategies:
        # Strategy box
        add_shape_box(slide, x_pos, 1.5, 3.5, 0.7,
                     f"{strategy['icon']} {strategy['name']}",
                     strategy['color'], WHITE, font_size=20, bold=True)

        # Description
        add_text_box(slide, x_pos, 2.3, 3.5, 0.6,
                    strategy['description'], font_size=13, color=DARK_GRAY)

        # Features
        add_text_box(slide, x_pos, 3, 3.5, 0.3,
                    "Key Features:", font_size=12, bold=True, color=strategy['color'])

        features_text = "\n".join([f"• {f}" for f in strategy['features']])
        add_text_box(slide, x_pos + 0.1, 3.35, 3.4, 1,
                    features_text, font_size=11, color=DARK_GRAY)

        # Use case box
        add_shape_box(slide, x_pos, 4.5, 3.5, 0.8,
                     f"Best For:\n{strategy['use_case']}",
                     LIGHT_GRAY, strategy['color'], font_size=11, bold=False)

        x_pos += 3.8

    # Decision matrix
    add_text_box(slide, 1.2, 5.5, 11, 0.3,
                "💡 Selection Guide: Choose based on risk tolerance, downtime requirements, and infrastructure capacity",
                font_size=14, color=DARK_BLUE, align=PP_ALIGN.CENTER)

def create_vm_architecture_overview_slide(prs):
    """Slide 6: VM Architecture Overview"""
    slide = add_slide_with_title(prs, "VM Deployment Architecture")

    y = 1.5

    # Harness Platform
    add_shape_box(slide, 3, y, 7, 0.7, "Harness CD Control Plane",
                 HARNESS_BLUE, WHITE, font_size=18)

    # Pipeline components
    y += 0.9
    comps = ["Pipeline", "Service", "Environment", "Infrastructure"]
    x = 3.5
    for comp in comps:
        add_shape_box(slide, x, y, 1.5, 0.5, comp, DARK_BLUE, WHITE, font_size=11)
        x += 1.6

    # Arrow
    add_arrow(slide, 6.5, y + 0.5, 6.5, y + 1, HARNESS_BLUE)

    # Delegate
    y += 1.2
    add_shape_box(slide, 4, y, 5, 0.6, "Harness Delegate\n(Customer Data Center/Private Cloud)",
                 WARNING_ORANGE, WHITE, font_size=14)

    add_text_box(slide, 5, y + 0.7, 3, 0.3, "SSH Connection", font_size=11, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)
    add_arrow(slide, 6.5, y + 0.6, 6.5, y + 1.1, WARNING_ORANGE)

    # VM Servers
    y += 1.4
    envs = [
        ("Dev\n2 VMs", 2),
        ("QA\n3 VMs", 3.5),
        ("Staging\n5 VMs", 5),
        ("Production\n10+ VMs", 6.5)
    ]

    for env, x_pos in envs:
        is_prod = "Production" in env
        color = ERROR_RED if is_prod else DARK_BLUE
        add_shape_box(slide, x_pos, y, 1.6, 0.7, env, color, WHITE, font_size=12)

    # Key points
    y += 1
    add_text_box(slide, 1, y, 11, 0.3, "Key Components:", font_size=16, bold=True, color=HARNESS_BLUE)

    points = [
        "🔌 SSH Connector: Secure connection using SSH keys or credentials",
        "📦 Artifact Sources: Docker Registry, Artifactory, Nexus, S3",
        "⚙️ Deployment Scripts: Pre-deploy, deploy, post-deploy, rollback",
        "✅ Health Checks: Process status, HTTP endpoints, smoke tests",
        "📊 Monitoring: Prometheus, Datadog, CloudWatch integration"
    ]

    y += 0.4
    for point in points:
        add_text_box(slide, 1.5, y, 10.5, 0.25, point, font_size=13, color=DARK_GRAY)
        y += 0.35

def create_vm_detailed_architecture_slide(prs):
    """Slide 7: VM Detailed Architecture"""
    slide = add_slide_with_title(prs, "VM Deployment: Detailed Architecture")

    # Left side - Components
    y = 1.5
    add_text_box(slide, 0.5, y, 5.5, 0.4, "Architecture Components", font_size=18, bold=True, color=HARNESS_BLUE)

    components = [
        ("Control Plane", [
            "Pipeline Orchestrator",
            "Service Definitions",
            "Environment Management",
            "Secrets Management",
            "RBAC & Governance"
        ]),
        ("Execution Layer", [
            "Harness Delegates (HA)",
            "SSH Connection Manager",
            "Artifact Download Engine",
            "Script Executor",
            "Health Check Monitor"
        ]),
        ("Target Infrastructure", [
            "VM Servers (Linux/Windows)",
            "Application Services",
            "Configuration Files",
            "Backup Storage",
            "Load Balancers"
        ])
    ]

    y += 0.5
    for comp_name, items in components:
        add_shape_box(slide, 0.5, y, 5.5, 0.4, comp_name, DARK_BLUE, WHITE, font_size=14)
        y += 0.5
        for item in items:
            add_text_box(slide, 0.7, y, 5.3, 0.25, f"• {item}", font_size=11, color=DARK_GRAY)
            y += 0.28
        y += 0.15

    # Right side - Deployment Flow
    x = 6.5
    y = 1.5
    add_text_box(slide, x, y, 6, 0.4, "Deployment Flow", font_size=18, bold=True, color=SUCCESS_GREEN)

    flow_steps = [
        ("1. Trigger", "Git webhook, manual, scheduled", HARNESS_BLUE),
        ("2. Artifact Fetch", "Download from registry", DARK_BLUE),
        ("3. Pre-Deploy", "Backup, stop services, validate", WARNING_ORANGE),
        ("4. Deploy", "Transfer, extract, configure", SUCCESS_GREEN),
        ("5. Post-Deploy", "Start services, health checks", SUCCESS_GREEN),
        ("6. Verify", "Continuous verification, metrics", HARNESS_BLUE),
        ("7. Complete", "Success or rollback", SUCCESS_GREEN)
    ]

    y += 0.5
    for step, desc, color in flow_steps:
        add_shape_box(slide, x, y, 2.5, 0.45, step, color, WHITE, font_size=13)
        add_text_box(slide, x + 2.6, y + 0.1, 3.3, 0.3, desc, font_size=11, color=DARK_GRAY)

        if step != "7. Complete":
            add_arrow(slide, x + 1.25, y + 0.45, x + 1.25, y + 0.65, color)

        y += 0.75

def create_vm_deployment_flow_phase1_slide(prs):
    """Slide 8: VM Deployment Flow - Phase 1 & 2"""
    slide = add_slide_with_title(prs, "VM Deployment Flow: Trigger & Pre-Deployment")

    # Phase 1: Trigger
    y = 1.5
    add_shape_box(slide, 1, y, 11, 0.5, "Phase 1: Trigger & Initialization",
                 HARNESS_BLUE, WHITE, font_size=18, bold=True)

    y += 0.6
    trigger_steps = [
        "1. Developer commits code to Git repository",
        "2. Webhook triggers Harness CD pipeline",
        "3. Pipeline fetches service configuration and artifact details",
        "4. Selects target environment (Dev/QA/Staging/Production)",
        "5. Assigns task to available Harness Delegate"
    ]
    add_bullet_list(slide, 1.5, y, 10, 1.2, trigger_steps, font_size=14)

    # Phase 2: Pre-Deployment
    y += 1.4
    add_shape_box(slide, 1, y, 11, 0.5, "Phase 2: Pre-Deployment Activities",
                 WARNING_ORANGE, WHITE, font_size=18, bold=True)

    y += 0.6

    # Create boxes for pre-deployment steps
    pre_deploy_boxes = [
        ("🔍 Environment\nValidation", "• Check disk space\n• Verify memory\n• Check dependencies"),
        ("💾 Backup\nCurrent App", "• Create backup dir\n• Copy files\n• Archive backup"),
        ("🛑 Stop\nServices", "• Stop application\n• Drain connections\n• Verify stopped")
    ]

    x = 1.5
    for title, details in pre_deploy_boxes:
        add_shape_box(slide, x, y, 3.2, 0.5, title, WARNING_ORANGE, WHITE, font_size=13, bold=True)
        add_text_box(slide, x, y + 0.55, 3.2, 0.8, details, font_size=10, color=DARK_GRAY)
        if x < 8:
            add_arrow(slide, x + 3.2, y + 0.25, x + 3.5, y + 0.25, WARNING_ORANGE)
        x += 3.5

    # Critical points
    y += 1.5
    add_text_box(slide, 1, y, 11, 0.3, "⚠️ Critical Safety Measures:", font_size=15, bold=True, color=ERROR_RED)
    y += 0.35
    safety = [
        "✓ Always create backup before stopping services - enables quick rollback if needed",
        "✓ Validate environment has sufficient resources - prevents deployment failures",
        "✓ Graceful service shutdown with connection draining - no dropped requests",
        "✓ Remove from load balancer pool before stopping - maintains availability"
    ]
    add_bullet_list(slide, 1.5, y, 10, 1, safety, font_size=13)

def create_vm_deployment_flow_phase2_slide(prs):
    """Slide 9: VM Deployment Flow - Deployment Execution"""
    slide = add_slide_with_title(prs, "VM Deployment Flow: Deployment Execution")

    y = 1.5
    add_shape_box(slide, 1, y, 11, 0.5, "Phase 3: Deployment Execution",
                 SUCCESS_GREEN, WHITE, font_size=18, bold=True)

    # Deployment steps as flow
    y += 0.7

    steps = [
        {
            "step": "1. Transfer Artifacts",
            "icon": "📦",
            "details": [
                "• SCP/SFTP transfer to target VM",
                "• Progress monitoring",
                "• Checksum verification"
            ],
            "command": "scp artifact.tar.gz user@vm:/tmp/"
        },
        {
            "step": "2. Extract & Install",
            "icon": "📂",
            "details": [
                "• Extract archive",
                "• Set permissions",
                "• Create directory structure"
            ],
            "command": "tar -xzf artifact.tar.gz -C /opt/myapp"
        },
        {
            "step": "3. Update Configuration",
            "icon": "⚙️",
            "details": [
                "• Copy environment configs",
                "• Replace variables",
                "• Inject secrets"
            ],
            "command": "cp config/prod/* /opt/myapp/conf/"
        },
        {
            "step": "4. Start Services",
            "icon": "▶️",
            "details": [
                "• systemctl start myapp",
                "• Verify process running",
                "• Check logs for errors"
            ],
            "command": "systemctl start myapp.service"
        }
    ]

    x = 1.2
    for i, step_info in enumerate(steps):
        # Step box
        add_shape_box(slide, x, y, 2.7, 0.4,
                     f"{step_info['icon']} {step_info['step']}",
                     SUCCESS_GREEN, WHITE, font_size=12, bold=True)

        # Details
        details_text = "\n".join(step_info['details'])
        add_text_box(slide, x, y + 0.45, 2.7, 0.8,
                    details_text, font_size=9, color=DARK_GRAY)

        # Command box
        add_shape_box(slide, x, y + 1.3, 2.7, 0.35,
                     step_info['command'],
                     LIGHT_GRAY, DARK_GRAY, font_size=8, bold=False)

        # Arrow to next step
        if i < len(steps) - 1:
            add_arrow(slide, x + 2.7, y + 0.2, x + 2.95, y + 0.2, SUCCESS_GREEN)

        x += 2.95

    # Rolling deployment note
    y += 1.8
    add_shape_box(slide, 2, y, 9, 0.8,
                 "🔄 Rolling Deployment Strategy\n" +
                 "Deploy to one VM at a time → Verify → Move to next VM\n" +
                 "Ensures continuous availability during deployment",
                 HARNESS_BLUE, WHITE, font_size=14, bold=False)

    # Timeline
    y += 1
    add_text_box(slide, 1, y, 11, 0.3,
                "⏱️ Typical Timeline: 5-7 minutes per VM | 10 VMs ≈ 60-70 minutes for complete deployment",
                font_size=13, color=DARK_BLUE, align=PP_ALIGN.CENTER)

def create_vm_deployment_flow_phase3_slide(prs):
    """Slide 10: VM Deployment Flow - Post-Deployment"""
    slide = add_slide_with_title(prs, "VM Deployment Flow: Post-Deployment & Verification")

    y = 1.5
    add_shape_box(slide, 1, y, 11, 0.5, "Phase 4: Post-Deployment Verification",
                 HARNESS_BLUE, WHITE, font_size=18, bold=True)

    y += 0.7

    # Verification steps in grid
    verifications = [
        {
            "title": "🏥 Health Checks",
            "checks": [
                "✓ Process status check",
                "✓ Port listening (8080)",
                "✓ HTTP /health endpoint",
                "✓ Database connectivity"
            ],
            "color": SUCCESS_GREEN
        },
        {
            "title": "🧪 Smoke Tests",
            "checks": [
                "✓ API endpoint tests",
                "✓ Authentication flow",
                "✓ Core functionality",
                "✓ External integrations"
            ],
            "color": HARNESS_BLUE
        },
        {
            "title": "⚖️ Load Balancer",
            "checks": [
                "✓ Register with LB pool",
                "✓ Wait for LB health checks",
                "✓ Verify traffic routing",
                "✓ Monitor connection count"
            ],
            "color": WARNING_ORANGE
        }
    ]

    x = 1.5
    for verify in verifications:
        add_shape_box(slide, x, y, 3.2, 0.4, verify['title'],
                     verify['color'], WHITE, font_size=14, bold=True)

        checks_text = "\n".join(verify['checks'])
        add_text_box(slide, x + 0.1, y + 0.45, 3, 1,
                    checks_text, font_size=11, color=DARK_GRAY)

        x += 3.5

    # Continuous Verification
    y += 1.6
    add_shape_box(slide, 1, y, 11, 0.5, "Phase 5: Continuous Verification (30 minutes)",
                 DARK_BLUE, WHITE, font_size=18, bold=True)

    y += 0.6

    # Metrics monitored
    add_text_box(slide, 1, y, 11, 0.3, "📊 Metrics Monitored:", font_size=16, bold=True, color=DARK_BLUE)

    y += 0.35
    metrics = [
        ("Error Rate", "< 1% (baseline comparison)", SUCCESS_GREEN),
        ("Response Time (P95)", "< baseline + 10%", SUCCESS_GREEN),
        ("CPU Usage", "< 80% sustained", SUCCESS_GREEN),
        ("Memory Usage", "< 85% sustained", SUCCESS_GREEN),
        ("Request Throughput", "≥ baseline - 5%", SUCCESS_GREEN)
    ]

    x = 1.5
    for metric, threshold, color in metrics:
        add_shape_box(slide, x, y, 2.1, 0.45,
                     f"{metric}\n{threshold}",
                     LIGHT_GRAY, color, font_size=10, bold=False)
        x += 2.2

    # Auto-rollback
    y += 0.6
    add_shape_box(slide, 2.5, y, 8, 0.7,
                 "🤖 AI-Powered Auto-Rollback\n" +
                 "Automatically reverts to previous version if anomalies detected\n" +
                 "Uses machine learning to compare against baseline metrics",
                 ERROR_RED, WHITE, font_size=13, bold=False)

def create_vm_deployment_pipeline_slide(prs):
    """Slide 11: VM Deployment Pipeline Configuration"""
    slide = add_slide_with_title(prs, "VM Deployment: Pipeline Configuration")

    y = 1.5

    # Pipeline structure
    add_text_box(slide, 1, y, 11, 0.3, "Pipeline Structure", font_size=18, bold=True, color=HARNESS_BLUE)

    y += 0.4

    # Service Definition
    add_shape_box(slide, 1, y, 5.5, 0.45, "📋 Service Definition",
                 HARNESS_BLUE, WHITE, font_size=14, bold=True)

    service_items = [
        "• Name: myapp-vm-service",
        "• Type: SSH",
        "• Artifacts:",
        "  - Docker: registry.example.com/myapp:v1.2.3",
        "  - TAR: artifactory.example.com/myapp.tar.gz",
        "• Config Files:",
        "  - Source: Git repository",
        "  - Paths: config/prod/*.properties",
        "• Variables:",
        "  - ENVIRONMENT: production",
        "  - LOG_LEVEL: INFO",
        "  - APP_PORT: 8080"
    ]
    add_bullet_list(slide, 1.1, y + 0.5, 5.3, 2.2, service_items, font_size=11)

    # Infrastructure Definition
    add_shape_box(slide, 6.7, y, 5.5, 0.45, "🏗️ Infrastructure Definition",
                 WARNING_ORANGE, WHITE, font_size=14, bold=True)

    infra_items = [
        "• Type: SSH",
        "• Connection:",
        "  - Host: prod-vm-*.example.com",
        "  - Port: 22",
        "  - Auth: SSH Key (from secrets)",
        "• Host Selection:",
        "  - Manual: Specific host list",
        "  - Dynamic: Tag-based discovery",
        "  - Load-balanced: Automatic distribution",
        "• Deployment:",
        "  - Strategy: Rolling update",
        "  - Batch size: 1 VM at a time",
        "  - Failure threshold: Stop on error"
    ]
    add_bullet_list(slide, 6.8, y + 0.5, 5.3, 2.2, infra_items, font_size=11)

    # Execution Steps
    y += 2.8
    add_shape_box(slide, 1, y, 11, 0.45, "⚙️ Execution Steps (Pipeline Stages)",
                 DARK_BLUE, WHITE, font_size=14, bold=True)

    stages = [
        ("Pre-Deploy", ["Backup", "Validate", "Stop"]),
        ("Deploy", ["Transfer", "Extract", "Configure"]),
        ("Post-Deploy", ["Start", "Health", "Tests"]),
        ("Verify", ["Metrics", "Logs", "Alerts"])
    ]

    x = 1.5
    for stage, steps in stages:
        add_shape_box(slide, x, y + 0.5, 2.4, 0.35, stage,
                     SUCCESS_GREEN, WHITE, font_size=12, bold=True)

        steps_text = "\n".join([f"• {s}" for s in steps])
        add_text_box(slide, x, y + 0.9, 2.4, 0.65,
                    steps_text, font_size=10, color=DARK_GRAY)

        if x < 9:
            add_arrow(slide, x + 2.4, y + 0.67, x + 2.6, y + 0.67, SUCCESS_GREEN)

        x += 2.6

    # Rollback configuration
    y += 1.7
    add_shape_box(slide, 2, y, 9, 0.55,
                 "🔄 Rollback Configuration: Automatic on failure | Manual trigger available | " +
                 "Restores from backup in < 5 minutes",
                 ERROR_RED, WHITE, font_size=12, bold=False)

def create_vm_rollback_slide(prs):
    """Slide 12: VM Rollback Strategy"""
    slide = add_slide_with_title(prs, "VM Deployment: Rollback Strategy")

    y = 1.5

    # Rollback triggers
    add_text_box(slide, 1, y, 11, 0.3, "🚨 Rollback Triggers", font_size=18, bold=True, color=ERROR_RED)

    y += 0.4
    triggers = [
        {
            "trigger": "Health Check Failure",
            "condition": "3 consecutive failures",
            "action": "Immediate rollback"
        },
        {
            "trigger": "Error Rate Spike",
            "condition": "> 5% for 5 minutes",
            "action": "Auto-rollback initiated"
        },
        {
            "trigger": "Performance Degradation",
            "condition": "Response time > 2x baseline",
            "action": "Alert & optional rollback"
        },
        {
            "trigger": "Manual Trigger",
            "condition": "Operator decision",
            "action": "Immediate rollback"
        }
    ]

    x = 1.2
    for trigger_info in triggers:
        add_shape_box(slide, x, y, 2.7, 0.35,
                     f"⚠️ {trigger_info['trigger']}",
                     ERROR_RED, WHITE, font_size=11, bold=True)

        add_text_box(slide, x, y + 0.4, 2.7, 0.5,
                    f"Condition:\n{trigger_info['condition']}\n\nAction:\n{trigger_info['action']}",
                    font_size=9, color=DARK_GRAY)

        x += 2.85

    # Rollback process
    y += 1.1
    add_text_box(slide, 1, y, 11, 0.3, "⏮️ Rollback Process", font_size=18, bold=True, color=DARK_BLUE)

    y += 0.4

    rollback_steps = [
        ("1. Detect Issue", "Monitoring alerts\nAuto or manual", WARNING_ORANGE),
        ("2. Stop Services", "Graceful shutdown", ERROR_RED),
        ("3. Restore Backup", "Latest good version", SUCCESS_GREEN),
        ("4. Start Services", "Verify health", SUCCESS_GREEN),
        ("5. Monitor", "Confirm stability", HARNESS_BLUE)
    ]

    x = 1.2
    for step, desc, color in rollback_steps:
        add_shape_box(slide, x, y, 2.1, 0.4, step, color, WHITE, font_size=11, bold=True)
        add_text_box(slide, x, y + 0.45, 2.1, 0.45, desc, font_size=10, color=DARK_GRAY)

        if x < 9.5:
            add_arrow(slide, x + 2.1, y + 0.2, x + 2.3, y + 0.2, color)

        x += 2.3

    # Best practices
    y += 1
    add_text_box(slide, 1, y, 11, 0.3, "✅ Rollback Best Practices", font_size=16, bold=True, color=SUCCESS_GREEN)

    y += 0.35
    practices = [
        "✓ Test rollback procedures monthly in non-production environments",
        "✓ Keep minimum 5 backup versions available at all times",
        "✓ Document rollback steps in runbooks for team reference",
        "✓ Set rollback time objective (RTO) < 5 minutes",
        "✓ Automate rollback process to reduce manual errors",
        "✓ Maintain separate backup storage location (not on deployment server)",
        "✓ Monitor post-rollback metrics for 30 minutes minimum"
    ]
    add_bullet_list(slide, 1.5, y, 10, 1.8, practices, font_size=13)

    # Rollback metrics
    y += 2
    add_shape_box(slide, 3, y, 7, 0.55,
                 "📊 Target Metrics: Mean Time To Detect (MTTD) < 2 min | " +
                 "Mean Time To Rollback (MTTR) < 5 min | Success Rate > 99%",
                 HARNESS_BLUE, WHITE, font_size=12, bold=False)

def create_ecs_architecture_overview_slide(prs):
    """Slide 13: ECS Architecture Overview"""
    slide = add_slide_with_title(prs, "ECS Deployment Architecture")

    y = 1.5

    # Harness Platform
    add_shape_box(slide, 3, y, 7, 0.7, "Harness CD Control Plane",
                 HARNESS_BLUE, WHITE, font_size=18)

    # Pipeline components
    y += 0.9
    comps = ["Pipeline", "Service", "Environment", "Infrastructure"]
    x = 3.5
    for comp in comps:
        add_shape_box(slide, x, y, 1.5, 0.5, comp, DARK_BLUE, WHITE, font_size=11)
        x += 1.6

    # Arrow
    add_arrow(slide, 6.5, y + 0.5, 6.5, y + 1, HARNESS_BLUE)

    # Delegate in AWS
    y += 1.2
    add_shape_box(slide, 4, y, 5, 0.6, "Harness Delegate\n(AWS VPC - Private Subnet)",
                 WARNING_ORANGE, WHITE, font_size=14)

    add_text_box(slide, 5, y + 0.7, 3, 0.3, "AWS API Calls", font_size=11, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)
    add_arrow(slide, 6.5, y + 0.6, 6.5, y + 1.1, WARNING_ORANGE)

    # AWS ECS Components
    y += 1.4

    # ECS Cluster
    add_shape_box(slide, 1.5, y, 3.3, 0.7, "ECS Cluster\n(Fargate/EC2)",
                 SUCCESS_GREEN, WHITE, font_size=13)

    # ECS Service
    add_shape_box(slide, 5.2, y, 3.3, 0.7, "ECS Service\n(10 tasks)",
                 SUCCESS_GREEN, WHITE, font_size=13)

    # Task Definition
    add_shape_box(slide, 8.9, y, 3.3, 0.7, "Task Definition\n(Container specs)",
                 SUCCESS_GREEN, WHITE, font_size=13)

    # Additional AWS services
    y += 1

    aws_services = [
        ("ALB\n(Load Balancer)", HARNESS_BLUE),
        ("ECR\n(Container Registry)", HARNESS_BLUE),
        ("CloudWatch\n(Monitoring)", HARNESS_BLUE)
    ]

    x = 2.5
    for service, color in aws_services:
        add_shape_box(slide, x, y, 2.5, 0.6, service, color, WHITE, font_size=12)
        x += 3

    # Key differences from VM
    y += 0.8
    add_text_box(slide, 1, y, 11, 0.3, "🔑 Key Differences from VM Deployment:", font_size=16, bold=True, color=HARNESS_BLUE)

    y += 0.35
    differences = [
        "✓ Container-based: Immutable infrastructure, consistent environments",
        "✓ AWS-managed: No server management, automatic scaling",
        "✓ Faster deployments: Parallel task updates, rapid scaling",
        "✓ Advanced strategies: Native Blue-Green and Canary support",
        "✓ Integration: Deep AWS ecosystem integration (ALB, ECR, CloudWatch)"
    ]
    add_bullet_list(slide, 1.5, y, 10, 1.3, differences, font_size=13)

def create_ecs_detailed_architecture_slide(prs):
    """Slide 14: ECS Detailed Architecture"""
    slide = add_slide_with_title(prs, "ECS Deployment: Detailed Architecture")

    # Left side - ECS Components
    y = 1.5
    add_text_box(slide, 0.5, y, 5.5, 0.4, "ECS Components", font_size=18, bold=True, color=HARNESS_BLUE)

    y += 0.5

    ecs_components = [
        ("ECS Cluster", SUCCESS_GREEN, [
            "Logical grouping of resources",
            "Fargate or EC2 launch type",
            "Multi-AZ for high availability",
            "Auto-scaling configuration"
        ]),
        ("ECS Service", SUCCESS_GREEN, [
            "Maintains desired task count",
            "Handles load balancer integration",
            "Manages deployments & rollbacks",
            "Health check configuration"
        ]),
        ("Task Definition", HARNESS_BLUE, [
            "Container specifications",
            "CPU & Memory limits",
            "Environment variables",
            "Volume mounts & secrets"
        ]),
        ("Tasks", DARK_BLUE, [
            "Running container instances",
            "Networked via VPC (awsvpc)",
            "Registered with target group",
            "Monitored via CloudWatch"
        ])
    ]

    for comp_name, color, items in ecs_components:
        add_shape_box(slide, 0.5, y, 5.5, 0.4, comp_name, color, WHITE, font_size=13)
        y += 0.45
        for item in items:
            add_text_box(slide, 0.7, y, 5.3, 0.2, f"• {item}", font_size=10, color=DARK_GRAY)
            y += 0.22
        y += 0.1

    # Right side - Deployment Flow
    x = 6.5
    y = 1.5
    add_text_box(slide, x, y, 6, 0.4, "ECS Deployment Flow", font_size=18, bold=True, color=SUCCESS_GREEN)

    y += 0.5

    flow = [
        ("1. Image Push to ECR", "CI pipeline builds & pushes", HARNESS_BLUE),
        ("2. Pipeline Trigger", "Webhook or manual", HARNESS_BLUE),
        ("3. Create Task Definition", "New revision with new image", WARNING_ORANGE),
        ("4. Update Service", "Rolling/Blue-Green/Canary", SUCCESS_GREEN),
        ("5. Start New Tasks", "Pull image, start containers", SUCCESS_GREEN),
        ("6. Register with ALB", "Health checks begin", SUCCESS_GREEN),
        ("7. Stop Old Tasks", "Graceful shutdown", WARNING_ORANGE),
        ("8. Verify Deployment", "Metrics & health checks", HARNESS_BLUE)
    ]

    for step, desc, color in flow:
        add_shape_box(slide, x, y, 2.8, 0.4, step, color, WHITE, font_size=12, bold=True)
        add_text_box(slide, x + 2.9, y + 0.08, 3, 0.25, desc, font_size=10, color=DARK_GRAY)

        if step != "8. Verify Deployment":
            add_arrow(slide, x + 1.4, y + 0.4, x + 1.4, y + 0.55, color)

        y += 0.62

def create_ecs_deployment_flow_phase1_slide(prs):
    """Slide 15: ECS Deployment Flow - Phase 1"""
    slide = add_slide_with_title(prs, "ECS Deployment Flow: Trigger & Pre-Deployment")

    y = 1.5

    # Phase 1: CI Pipeline & ECR
    add_shape_box(slide, 1, y, 11, 0.5, "Phase 1: CI Pipeline & Container Registry",
                 HARNESS_BLUE, WHITE, font_size=18, bold=True)

    y += 0.6

    ci_steps = [
        ("Build", "Docker build\n-t myapp:v1.2.3", HARNESS_BLUE),
        ("Test", "Run unit tests\nSecurity scan", SUCCESS_GREEN),
        ("Push", "docker push ECR\nmyapp:v1.2.3", WARNING_ORANGE),
        ("Trigger", "SNS/Webhook\n→ Harness", ERROR_RED)
    ]

    x = 1.8
    for step, desc, color in ci_steps:
        add_shape_box(slide, x, y, 2.2, 0.5, step, color, WHITE, font_size=14, bold=True)
        add_text_box(slide, x, y + 0.55, 2.2, 0.5, desc, font_size=10, color=DARK_GRAY)

        if x < 9:
            add_arrow(slide, x + 2.2, y + 0.25, x + 2.4, y + 0.25, color)

        x += 2.4

    # Phase 2: Pre-Deployment
    y += 1.2
    add_shape_box(slide, 1, y, 11, 0.5, "Phase 2: Pre-Deployment Validation",
                 WARNING_ORANGE, WHITE, font_size=18, bold=True)

    y += 0.6

    validation_steps = [
        {
            "title": "🔐 AWS Authentication",
            "items": [
                "• Assume IAM role",
                "• Verify ECS permissions",
                "• Check ECR access",
                "• Validate ALB permissions"
            ]
        },
        {
            "title": "🖼️ Image Verification",
            "items": [
                "• Confirm image exists in ECR",
                "• Verify image digest",
                "• Check vulnerability scan",
                "• Validate image size"
            ]
        },
        {
            "title": "📊 Capacity Check",
            "items": [
                "• Check Fargate quotas",
                "• Verify subnet capacity",
                "• Validate security groups",
                "• Check target group health"
            ]
        }
    ]

    x = 1.2
    for validation in validation_steps:
        add_shape_box(slide, x, y, 3.5, 0.4, validation['title'],
                     WARNING_ORANGE, WHITE, font_size=13, bold=True)

        items_text = "\n".join(validation['items'])
        add_text_box(slide, x + 0.1, y + 0.45, 3.3, 1,
                    items_text, font_size=10, color=DARK_GRAY)

        x += 3.7

    # Current state check
    y += 1.6
    add_text_box(slide, 1, y, 11, 0.3, "📋 Current State Analysis:", font_size=16, bold=True, color=DARK_BLUE)

    y += 0.35

    current_state = [
        "✓ Retrieve current service configuration (desired count, task definition)",
        "✓ Get current task definition revision and container image version",
        "✓ Check running task count and health status",
        "✓ Verify ALB target group registration and health checks",
        "✓ Capture baseline metrics for comparison (error rate, response time, CPU)"
    ]
    add_bullet_list(slide, 1.5, y, 10, 1.2, current_state, font_size=13)

    # Save for rollback note
    y += 1.4
    add_shape_box(slide, 2.5, y, 8, 0.5,
                 "💾 Current state saved as rollback point - enables instant revert if needed",
                 DARK_BLUE, WHITE, font_size=13, bold=False)

def create_ecs_deployment_flow_phase2_slide(prs):
    """Slide 16: ECS Deployment Flow - Phase 2"""
    slide = add_slide_with_title(prs, "ECS Deployment Flow: Task Definition & Service Update")

    y = 1.5

    # Phase 3: Task Definition
    add_shape_box(slide, 1, y, 11, 0.5, "Phase 3: Task Definition Creation",
                 HARNESS_BLUE, WHITE, font_size=18, bold=True)

    y += 0.6

    # Task definition components
    add_text_box(slide, 1, y, 5.5, 0.3, "Task Definition Structure:", font_size=14, bold=True, color=HARNESS_BLUE)

    task_def_items = [
        "• Family: myapp-task",
        "• Revision: 15 (auto-incremented)",
        "• Container Definition:",
        "  - Name: myapp-container",
        "  - Image: ECR/myapp:v1.2.3",
        "  - CPU: 256 (0.25 vCPU)",
        "  - Memory: 512 MB",
        "  - Port Mappings: 8080:8080",
        "  - Environment Variables",
        "  - Secrets from AWS Secrets Manager",
        "• Network Mode: awsvpc",
        "• Launch Type: FARGATE",
        "• IAM Roles: Task & Execution roles"
    ]
    add_bullet_list(slide, 1.1, y + 0.35, 5.3, 2.3, task_def_items, font_size=11)

    # Register task definition
    x = 6.7
    add_text_box(slide, x, y, 5.5, 0.3, "Register New Revision:", font_size=14, bold=True, color=SUCCESS_GREEN)

    register_steps = [
        "1. Clone current task definition",
        "2. Update container image:",
        "   FROM: myapp:v1.2.2",
        "   TO: myapp:v1.2.3",
        "3. Update environment variables",
        "4. Inject new secrets if needed",
        "5. Validate task definition",
        "6. Register with AWS ECS:",
        "   aws ecs register-task-definition",
        "7. Receive new ARN:",
        "   arn:aws:ecs:...:task-def/",
        "   myapp-task:15",
        "8. Ready for deployment ✓"
    ]
    add_bullet_list(slide, x + 0.1, y + 0.35, 5.3, 2.3, register_steps, font_size=11)

    # Phase 4: Service Update
    y += 2.5
    add_shape_box(slide, 1, y, 11, 0.5, "Phase 4: ECS Service Update",
                 SUCCESS_GREEN, WHITE, font_size=18, bold=True)

    y += 0.6

    # Deployment configuration
    deployment_config = [
        {
            "param": "Deployment Strategy",
            "value": "Rolling / Blue-Green / Canary",
            "color": HARNESS_BLUE
        },
        {
            "param": "Maximum Percent",
            "value": "200% (can run 20 tasks during deployment)",
            "color": SUCCESS_GREEN
        },
        {
            "param": "Minimum Healthy",
            "value": "100% (maintain 10 healthy tasks minimum)",
            "color": SUCCESS_GREEN
        },
        {
            "param": "Circuit Breaker",
            "value": "Enabled - Auto-rollback on failures",
            "color": WARNING_ORANGE
        }
    ]

    for config in deployment_config:
        add_shape_box(slide, 1.5, y, 3, 0.35, config['param'],
                     config['color'], WHITE, font_size=12, bold=True)
        add_text_box(slide, 4.6, y + 0.08, 6.8, 0.25, config['value'],
                    font_size=11, color=DARK_GRAY)
        y += 0.45

    # Update command
    y += 0.2
    add_shape_box(slide, 2, y, 9, 0.6,
                 "aws ecs update-service \\\n" +
                 "  --cluster prod-ecs-cluster \\\n" +
                 "  --service myapp-service \\\n" +
                 "  --task-definition myapp-task:15 \\\n" +
                 "  --force-new-deployment",
                 LIGHT_GRAY, DARK_GRAY, font_size=11, bold=False)

def create_ecs_blue_green_detailed_slide(prs):
    """Slide 17: ECS Blue-Green Deployment Detailed"""
    slide = add_slide_with_title(prs, "ECS Blue-Green Deployment Strategy")

    y = 1.5

    # Strategy overview
    add_shape_box(slide, 1, y, 11, 0.5,
                 "🔵🟢 Blue-Green: Zero-Downtime Deployment with Instant Rollback",
                 SUCCESS_GREEN, WHITE, font_size=18, bold=True)

    y += 0.6

    # Visual representation of stages
    stages = [
        {
            "stage": "Stage 1: Initial",
            "blue": "100% Traffic",
            "blue_color": HARNESS_BLUE,
            "green": "0% Traffic",
            "green_color": LIGHT_GRAY,
            "desc": "Blue serving all production traffic"
        },
        {
            "stage": "Stage 2: Deploy Green",
            "blue": "100% Traffic",
            "blue_color": HARNESS_BLUE,
            "green": "0% (Testing)",
            "green_color": SUCCESS_GREEN,
            "desc": "Green deployed, not receiving traffic"
        },
        {
            "stage": "Stage 3: Shift 10%",
            "blue": "90% Traffic",
            "blue_color": HARNESS_BLUE,
            "green": "10% Traffic",
            "green_color": SUCCESS_GREEN,
            "desc": "Initial traffic validation (5 min)"
        },
        {
            "stage": "Stage 4: Shift 50%",
            "blue": "50% Traffic",
            "blue_color": HARNESS_BLUE,
            "green": "50% Traffic",
            "green_color": SUCCESS_GREEN,
            "desc": "Half traffic validation (10 min)"
        },
        {
            "stage": "Stage 5: Complete",
            "blue": "0% (Standby)",
            "blue_color": LIGHT_GRAY,
            "green": "100% Traffic",
            "green_color": SUCCESS_GREEN,
            "desc": "Green serving all traffic, blue ready for rollback"
        }
    ]

    x = 0.8
    for stage in stages:
        # Stage title
        add_text_box(slide, x, y, 2.2, 0.25, stage['stage'], font_size=10, bold=True, color=DARK_BLUE)

        # Blue environment
        add_shape_box(slide, x, y + 0.3, 2.2, 0.4, f"Blue\n{stage['blue']}",
                     stage['blue_color'], WHITE, font_size=10, bold=True)

        # Green environment
        add_shape_box(slide, x, y + 0.75, 2.2, 0.4, f"Green\n{stage['green']}",
                     stage['green_color'], WHITE if stage['green_color'] != LIGHT_GRAY else DARK_GRAY,
                     font_size=10, bold=True)

        # Description
        add_text_box(slide, x, y + 1.2, 2.2, 0.4, stage['desc'], font_size=8, color=DARK_GRAY)

        # Arrow to next stage
        if x < 9:
            add_arrow(slide, x + 2.2, y + 0.8, x + 2.4, y + 0.8, SUCCESS_GREEN)

        x += 2.4

    # ALB Configuration
    y += 1.8
    add_text_box(slide, 1, y, 11, 0.3, "⚖️ Application Load Balancer Configuration:", font_size=16, bold=True, color=HARNESS_BLUE)

    y += 0.35

    alb_config = [
        {
            "component": "Blue Target Group",
            "details": "Original tasks, gradually reduced traffic"
        },
        {
            "component": "Green Target Group",
            "details": "New tasks, gradually increased traffic"
        },
        {
            "component": "Listener Rules",
            "details": "Weight-based routing: Blue:90 → Green:10 → ..."
        },
        {
            "component": "Health Checks",
            "details": "/health endpoint, 30s interval, 2/3 threshold"
        }
    ]

    x = 1.5
    for config in alb_config:
        add_shape_box(slide, x, y, 2.4, 0.35, config['component'],
                     HARNESS_BLUE, WHITE, font_size=11, bold=True)
        add_text_box(slide, x, y + 0.4, 2.4, 0.4, config['details'], font_size=9, color=DARK_GRAY)
        x += 2.6

    # Benefits
    y += 0.9
    add_text_box(slide, 1, y, 11, 0.3, "✅ Key Benefits:", font_size=16, bold=True, color=SUCCESS_GREEN)

    y += 0.35
    benefits = [
        "✓ Zero downtime: Traffic switch is instantaneous at ALB level",
        "✓ Full validation: Complete testing before receiving production traffic",
        "✓ Instant rollback: Simply redirect traffic back to blue (< 1 minute)",
        "✓ Safe deployment: Any issues detected only affect small percentage initially"
    ]
    add_bullet_list(slide, 1.5, y, 10, 0.9, benefits, font_size=13)

    # Timeline
    y += 1
    add_shape_box(slide, 2.5, y, 8, 0.45,
                 "⏱️ Total Deployment Time: ~45-55 minutes (including 30-min verification)",
                 DARK_BLUE, WHITE, font_size=13, bold=False)

def create_ecs_canary_detailed_slide(prs):
    """Slide 18: ECS Canary Deployment Detailed"""
    slide = add_slide_with_title(prs, "ECS Canary Deployment Strategy")

    y = 1.5

    # Strategy overview
    add_shape_box(slide, 1, y, 11, 0.5,
                 "🐤 Canary: Gradual Rollout with Progressive Validation",
                 WARNING_ORANGE, WHITE, font_size=18, bold=True)

    y += 0.6

    # Canary phases
    add_text_box(slide, 1, y, 11, 0.3, "Canary Deployment Phases:", font_size=16, bold=True, color=WARNING_ORANGE)

    y += 0.4

    canary_phases = [
        {
            "phase": "Phase 1: 10% Canary",
            "tasks": "1 new task, 9 old tasks",
            "duration": "5 minutes",
            "validation": "Error rate < 1%, Response time baseline",
            "color": WARNING_ORANGE
        },
        {
            "phase": "Phase 2: 25% Canary",
            "tasks": "2-3 new tasks, 7-8 old tasks",
            "duration": "10 minutes",
            "validation": "Extended metrics, CPU/Memory checks",
            "color": WARNING_ORANGE
        },
        {
            "phase": "Phase 3: 50% Canary",
            "tasks": "5 new tasks, 5 old tasks",
            "duration": "15 minutes",
            "validation": "Full metrics suite, business KPIs",
            "color": SUCCESS_GREEN
        },
        {
            "phase": "Phase 4: 100% Rollout",
            "tasks": "All 10 new tasks",
            "duration": "Ongoing",
            "validation": "Continuous monitoring, auto-rollback ready",
            "color": SUCCESS_GREEN
        }
    ]

    for i, phase in enumerate(canary_phases):
        # Phase box
        add_shape_box(slide, 1, y, 11, 0.35, phase['phase'],
                     phase['color'], WHITE, font_size=13, bold=True)

        y += 0.4

        # Details
        details = [
            f"Tasks: {phase['tasks']}",
            f"Duration: {phase['duration']}",
            f"Validation: {phase['validation']}"
        ]

        x = 1.5
        for detail in details:
            add_text_box(slide, x, y, 3.3, 0.2, f"• {detail}", font_size=10, color=DARK_GRAY)
            x += 3.5

        # Arrow to next phase
        if i < len(canary_phases) - 1:
            add_arrow(slide, 6.5, y + 0.25, 6.5, y + 0.45, phase['color'])

        y += 0.5

    # Monitoring & Auto-Rollback
    y += 0.2
    add_text_box(slide, 1, y, 11, 0.3, "🤖 Automated Validation & Rollback:", font_size=16, bold=True, color=DARK_BLUE)

    y += 0.35

    validation_metrics = [
        {
            "metric": "Error Rate",
            "threshold": "< 1%",
            "action": "Auto-rollback if exceeded"
        },
        {
            "metric": "Response Time",
            "threshold": "< Baseline + 10%",
            "action": "Alert if exceeded, rollback if > 50%"
        },
        {
            "metric": "HTTP 5xx Errors",
            "threshold": "< 10 per minute",
            "action": "Immediate rollback"
        },
        {
            "metric": "Resource Usage",
            "threshold": "CPU < 80%, Memory < 85%",
            "action": "Alert, continue monitoring"
        }
    ]

    for metric in validation_metrics:
        add_shape_box(slide, 1.5, y, 2.5, 0.3, metric['metric'],
                     HARNESS_BLUE, WHITE, font_size=11, bold=True)
        add_text_box(slide, 4.1, y + 0.05, 2.5, 0.22, metric['threshold'], font_size=10, color=DARK_GRAY)
        add_text_box(slide, 6.7, y + 0.05, 4.5, 0.22, metric['action'], font_size=10, color=ERROR_RED)
        y += 0.38

    # Benefits
    y += 0.1
    add_text_box(slide, 1, y, 11, 0.3, "✅ Why Use Canary?", font_size=16, bold=True, color=SUCCESS_GREEN)

    y += 0.35
    benefits = [
        "✓ Risk mitigation: Limits blast radius to small percentage of traffic initially",
        "✓ Progressive validation: Extended validation time at each stage",
        "✓ Automatic safety: AI-powered rollback on any metric degradation",
        "✓ Business-safe: Minimal impact if issues occur (only affects canary percentage)"
    ]
    add_bullet_list(slide, 1.5, y, 10, 0.9, benefits, font_size=13)

def create_ecs_pipeline_slide(prs):
    """Slide 19: ECS Pipeline Configuration"""
    slide = add_slide_with_title(prs, "ECS Deployment: Pipeline Configuration")

    y = 1.5

    # Service Definition
    add_shape_box(slide, 1, y, 5.5, 0.45, "📋 Service Definition",
                 HARNESS_BLUE, WHITE, font_size=14, bold=True)

    service_items = [
        "Service: myapp-ecs-service",
        "Type: ECS",
        "",
        "Artifact:",
        "• Type: ECR",
        "• Repository: myapp",
        "• Registry: 123...amazonaws.com",
        "• Tag: <+pipeline.sequenceId>",
        "",
        "Task Definition:",
        "• Family: myapp-task",
        "• CPU: 256 (.25 vCPU)",
        "• Memory: 512 MB",
        "• Network Mode: awsvpc",
        "• Launch Type: FARGATE"
    ]
    add_bullet_list(slide, 1.1, y + 0.5, 5.3, 2.8, service_items, font_size=11)

    # Infrastructure Definition
    add_shape_box(slide, 6.7, y, 5.5, 0.45, "🏗️ Infrastructure Definition",
                 WARNING_ORANGE, WHITE, font_size=14, bold=True)

    infra_items = [
        "Type: ECS",
        "AWS Connector: aws-prod",
        "",
        "ECS Configuration:",
        "• Cluster: prod-ecs-cluster",
        "• Region: us-east-1",
        "• Launch Type: FARGATE",
        "",
        "Network Configuration:",
        "• Subnets: private-1a, private-1b, private-1c",
        "• Security Groups: ecs-task-sg",
        "• Assign Public IP: DISABLED",
        "",
        "Load Balancer:",
        "• ALB: myapp-alb",
        "• Target Group: myapp-tg",
        "• Container Port: 8080"
    ]
    add_bullet_list(slide, 6.8, y + 0.5, 5.3, 2.8, infra_items, font_size=11)

    # Deployment Configuration
    y += 3.4
    add_shape_box(slide, 1, y, 11, 0.45, "⚙️ Deployment Configuration",
                 DARK_BLUE, WHITE, font_size=14, bold=True)

    y += 0.5

    deployment_configs = [
        ("Strategy", ["Rolling Update", "Blue-Green", "Canary"], SUCCESS_GREEN),
        ("Health Check", ["Grace Period: 60s", "Interval: 30s", "Path: /health"], HARNESS_BLUE),
        ("Scaling", ["Desired: 10 tasks", "Max: 200%", "Min: 100%"], WARNING_ORANGE),
        ("Circuit Breaker", ["Enabled: true", "Rollback: automatic", "Failure %: 10%"], ERROR_RED)
    ]

    x = 1.2
    for config_type, configs, color in deployment_configs:
        add_shape_box(slide, x, y, 2.5, 0.35, config_type, color, WHITE, font_size=12, bold=True)

        config_text = "\n".join([f"• {c}" for c in configs])
        add_text_box(slide, x, y + 0.4, 2.5, 0.7, config_text, font_size=10, color=DARK_GRAY)

        x += 2.7

    # Pipeline Stages
    y += 1.2
    add_text_box(slide, 1, y, 11, 0.3, "Pipeline Stages:", font_size=14, bold=True, color=HARNESS_BLUE)

    y += 0.35
    stages_text = "Validate → Create Task Definition → Update Service → Wait for Steady State → Health Checks → Verify Metrics → Complete"
    add_text_box(slide, 1.5, y, 10, 0.5, stages_text, font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

def create_network_architecture_slide(prs):
    """Slide 20: Network Architecture"""
    slide = add_slide_with_title(prs, "Network Architecture")

    y = 1.5

    # VPC Overview
    add_shape_box(slide, 1, y, 11, 0.5, "AWS VPC Architecture (Multi-AZ for High Availability)",
                 HARNESS_BLUE, WHITE, font_size=16, bold=True)

    y += 0.6

    # Availability Zones
    azs = [
        ("AZ-1a", 1.5),
        ("AZ-1b", 5),
        ("AZ-1c", 8.5)
    ]

    for az_name, x_start in azs:
        # AZ Label
        add_text_box(slide, x_start, y, 3, 0.25, az_name, font_size=12, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

        # Public Subnet
        add_shape_box(slide, x_start, y + 0.3, 3, 0.5, "Public Subnet\nALB, NAT Gateway",
                     SUCCESS_GREEN, WHITE, font_size=10, bold=True)

        # Private Subnet (ECS)
        add_shape_box(slide, x_start, y + 0.85, 3, 0.5, "Private Subnet\nECS Tasks, VMs",
                     HARNESS_BLUE, WHITE, font_size=10, bold=True)

        # Private Subnet (Delegate)
        add_shape_box(slide, x_start, y + 1.4, 3, 0.5, "Private Subnet\nHarness Delegates",
                     WARNING_ORANGE, WHITE, font_size=10, bold=True)

    # Network Components
    y += 2
    add_text_box(slide, 1, y, 11, 0.3, "Network Components:", font_size=16, bold=True, color=HARNESS_BLUE)

    y += 0.35

    network_components = [
        "🌐 Internet Gateway: Public internet access for ALB",
        "🔄 NAT Gateways: Outbound internet for private subnets (one per AZ)",
        "⚖️ Application Load Balancer: Distributes traffic to ECS tasks/VMs",
        "🔗 VPC Endpoints: Private AWS service access (ECS, ECR, S3, CloudWatch)",
        "🔒 Security Groups: Stateful firewall rules for resources",
        "📊 VPC Flow Logs: Network traffic logging to CloudWatch"
    ]
    add_bullet_list(slide, 1.5, y, 10, 1.5, network_components, font_size=13)

    # Security Groups
    y += 1.6
    add_text_box(slide, 1, y, 11, 0.3, "Security Group Rules:", font_size=16, bold=True, color=ERROR_RED)

    y += 0.35

    sg_rules = [
        {
            "sg": "ALB SG",
            "inbound": "80, 443 from 0.0.0.0/0",
            "outbound": "8080 to ECS Task SG"
        },
        {
            "sg": "ECS Task SG",
            "inbound": "8080 from ALB SG",
            "outbound": "443 to 0.0.0.0/0, 5432 to RDS SG"
        },
        {
            "sg": "Delegate SG",
            "inbound": "None",
            "outbound": "443 to Harness, 22 to VM SG, AWS APIs"
        }
    ]

    for rule in sg_rules:
        add_shape_box(slide, 1.5, y, 2, 0.3, rule['sg'], DARK_BLUE, WHITE, font_size=11, bold=True)
        add_text_box(slide, 3.6, y + 0.05, 3.5, 0.22, f"In: {rule['inbound']}", font_size=9, color=DARK_GRAY)
        add_text_box(slide, 7.2, y + 0.05, 4, 0.22, f"Out: {rule['outbound']}", font_size=9, color=DARK_GRAY)
        y += 0.38

    # Network Traffic Flow
    y += 0.1
    add_shape_box(slide, 2, y, 9, 0.5,
                 "User → Internet Gateway → ALB (Public) → ECS Tasks/VMs (Private) → NAT Gateway → AWS Services",
                 SUCCESS_GREEN, WHITE, font_size=11, bold=False)

def create_security_architecture_slide(prs):
    """Slide 21: Security Architecture"""
    slide = add_slide_with_title(prs, "Security Architecture")

    y = 1.5

    # Security Layers
    security_layers = [
        {
            "layer": "1. Network Security",
            "icon": "🔒",
            "color": DARK_BLUE,
            "items": [
                "VPC isolation with private subnets",
                "Security groups (stateful firewall)",
                "Network ACLs (stateless firewall)",
                "No inbound connections to delegates"
            ]
        },
        {
            "layer": "2. Access Control",
            "icon": "👥",
            "color": HARNESS_BLUE,
            "items": [
                "Role-Based Access Control (RBAC)",
                "IAM roles for AWS resources",
                "SSH key authentication (VMs)",
                "MFA for Harness platform access"
            ]
        },
        {
            "layer": "3. Secrets Management",
            "icon": "🔐",
            "color": WARNING_ORANGE,
            "items": [
                "AWS Secrets Manager integration",
                "HashiCorp Vault support",
                "Encrypted at rest & in transit",
                "Automatic secret rotation"
            ]
        },
        {
            "layer": "4. Encryption",
            "icon": "🛡️",
            "color": SUCCESS_GREEN,
            "items": [
                "TLS 1.3 for all communication",
                "mTLS for delegate connections",
                "Encrypted storage (EBS, RDS)",
                "SSL/TLS termination at ALB"
            ]
        }
    ]

    row = 0
    for i, layer in enumerate(security_layers):
        x = 1.5 if i % 2 == 0 else 6.7
        y_pos = y + (row * 1.5)

        add_shape_box(slide, x, y_pos, 4.8, 0.4, f"{layer['icon']} {layer['layer']}",
                     layer['color'], WHITE, font_size=13, bold=True)

        items_text = "\n".join([f"• {item}" for item in layer['items']])
        add_text_box(slide, x + 0.1, y_pos + 0.45, 4.6, 0.9,
                    items_text, font_size=10, color=DARK_GRAY)

        if i % 2 == 1:
            row += 1

    # Compliance & Audit
    y += 3.2
    add_text_box(slide, 1, y, 11, 0.3, "Compliance & Audit:", font_size=16, bold=True, color=HARNESS_BLUE)

    y += 0.35

    compliance_items = [
        "📋 Complete Audit Trails: Every action logged with timestamp, user, and change details",
        "✅ Compliance Ready: SOC 2 Type II, ISO 27001, GDPR, HIPAA compliant",
        "🔍 Policy as Code: OPA integration for automated policy enforcement",
        "📊 Compliance Reports: Automated generation for audits",
        "🚨 Alert on Policy Violations: Real-time notifications for policy breaches"
    ]
    add_bullet_list(slide, 1.5, y, 10, 1.2, compliance_items, font_size=13)

    # Security best practices
    y += 1.3
    add_shape_box(slide, 2, y, 9, 0.7,
                 "🔐 Security Best Practices\n" +
                 "Least privilege access | Regular key rotation | Security scanning | " +
                 "Vulnerability management | Incident response plan",
                 ERROR_RED, WHITE, font_size=12, bold=False)

def create_monitoring_slide(prs):
    """Slide 22: Monitoring & Verification"""
    slide = add_slide_with_title(prs, "Monitoring & Continuous Verification")

    y = 1.5

    # Monitoring Tools
    add_text_box(slide, 1, y, 11, 0.3, "Monitoring Stack:", font_size=16, bold=True, color=HARNESS_BLUE)

    y += 0.4

    monitoring_tools = [
        {
            "tool": "Prometheus",
            "logo": "📊",
            "use_case": "Metrics Collection",
            "features": ["Time-series database", "PromQL queries", "Alertmanager integration"],
            "color": WARNING_ORANGE
        },
        {
            "tool": "Datadog",
            "logo": "🐕",
            "use_case": "APM & Monitoring",
            "features": ["Application Performance", "Infrastructure monitoring", "Log aggregation"],
            "color": HARNESS_BLUE
        },
        {
            "tool": "CloudWatch",
            "logo": "☁️",
            "use_case": "AWS Native",
            "features": ["ECS Container Insights", "ALB metrics", "Custom metrics & alarms"],
            "color": SUCCESS_GREEN
        }
    ]

    x = 1.5
    for tool in monitoring_tools:
        add_shape_box(slide, x, y, 3.2, 0.4, f"{tool['logo']} {tool['tool']}",
                     tool['color'], WHITE, font_size=14, bold=True)

        add_text_box(slide, x, y + 0.45, 3.2, 0.25, tool['use_case'], font_size=11, bold=True, color=tool['color'])

        features_text = "\n".join([f"• {f}" for f in tool['features']])
        add_text_box(slide, x + 0.1, y + 0.72, 3, 0.7, features_text, font_size=9, color=DARK_GRAY)

        x += 3.5

    # Continuous Verification
    y += 1.5
    add_shape_box(slide, 1, y, 11, 0.5, "🤖 Harness Continuous Verification (AI-Powered)",
                 DARK_BLUE, WHITE, font_size=16, bold=True)

    y += 0.6

    cv_features = [
        {
            "feature": "Baseline Comparison",
            "description": "Compares current deployment metrics against previous successful deployments"
        },
        {
            "feature": "Anomaly Detection",
            "description": "Machine learning identifies unusual patterns in metrics and logs"
        },
        {
            "feature": "Multi-Source Analysis",
            "description": "Aggregates data from Prometheus, Datadog, CloudWatch, and custom sources"
        },
        {
            "feature": "Auto-Rollback Decision",
            "description": "Automatically triggers rollback if metrics exceed thresholds or anomalies detected"
        }
    ]

    for feature in cv_features:
        add_shape_box(slide, 1.5, y, 3.5, 0.35, feature['feature'],
                     HARNESS_BLUE, WHITE, font_size=12, bold=True)
        add_text_box(slide, 5.1, y + 0.05, 6, 0.25, feature['description'], font_size=11, color=DARK_GRAY)
        y += 0.45

    # Key Metrics
    y += 0.2
    add_text_box(slide, 1, y, 11, 0.3, "Key Metrics Monitored:", font_size=16, bold=True, color=SUCCESS_GREEN)

    y += 0.35

    key_metrics = [
        "📉 Error Rates: HTTP 4xx/5xx errors, application exceptions",
        "⏱️ Response Times: P50, P95, P99 latency measurements",
        "💻 Resource Usage: CPU, memory, disk I/O, network bandwidth",
        "📊 Throughput: Requests per second, transactions per minute",
        "🎯 Business Metrics: Conversion rates, user actions, custom KPIs",
        "🔗 Dependencies: Database latency, external API response times"
    ]
    add_bullet_list(slide, 1.5, y, 10, 1.5, key_metrics, font_size=13)

    # Verification timeline
    y += 1.6
    add_shape_box(slide, 2.5, y, 8, 0.5,
                 "⏱️ Verification Timeline: 5-minute canary → 10-minute blue-green → 30-minute production soak",
                 WARNING_ORANGE, WHITE, font_size=12, bold=False)

def create_deployment_comparison_slide(prs):
    """Slide 23: Deployment Strategy Comparison"""
    slide = add_slide_with_title(prs, "Deployment Strategy Comparison")

    y = 1.5

    # Comparison table header
    headers = ["Strategy", "Downtime", "Rollback Speed", "Resource Cost", "Complexity", "Best For"]
    x_positions = [1, 3, 4.8, 6.6, 8.4, 10]
    widths = [1.8, 1.6, 1.6, 1.6, 1.6, 2]

    for header, x, width in zip(headers, x_positions, widths):
        add_shape_box(slide, x, y, width, 0.4, header, DARK_BLUE, WHITE, font_size=12, bold=True)

    # Comparison data
    comparisons = [
        {
            "strategy": "Rolling Update",
            "downtime": "Minimal",
            "rollback": "Manual\n5-10 min",
            "cost": "Low\n(1x)",
            "complexity": "Low",
            "best_for": "Standard\ndeployments",
            "color": HARNESS_BLUE
        },
        {
            "strategy": "Blue-Green",
            "downtime": "Zero",
            "rollback": "Instant\n< 1 min",
            "cost": "High\n(2x)",
            "complexity": "Medium",
            "best_for": "Critical\napplications",
            "color": SUCCESS_GREEN
        },
        {
            "strategy": "Canary",
            "downtime": "Zero",
            "rollback": "Auto\n2-3 min",
            "cost": "Medium\n(1.1-1.5x)",
            "complexity": "High",
            "best_for": "High-risk\nchanges",
            "color": WARNING_ORANGE
        }
    ]

    y += 0.5
    for comparison in comparisons:
        # Strategy name
        add_shape_box(slide, 1, y, 1.8, 0.45, comparison['strategy'],
                     comparison['color'], WHITE, font_size=11, bold=True)

        # Comparison values
        values = [
            comparison['downtime'],
            comparison['rollback'],
            comparison['cost'],
            comparison['complexity'],
            comparison['best_for']
        ]

        x = 3
        for value, width in zip(values, widths[1:]):
            add_text_box(slide, x, y + 0.08, width, 0.35, value, font_size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)
            x += width + 0.2

        y += 0.55

    # Decision Guide
    y += 0.3
    add_text_box(slide, 1, y, 11, 0.3, "Decision Guide:", font_size=16, bold=True, color=HARNESS_BLUE)

    y += 0.35

    decision_guide = [
        {
            "scenario": "Low-risk, frequent deployments",
            "recommendation": "Rolling Update",
            "rationale": "Cost-effective, simple, sufficient for most cases"
        },
        {
            "scenario": "Zero-downtime requirement",
            "recommendation": "Blue-Green",
            "rationale": "Instant switch, full validation, immediate rollback"
        },
        {
            "scenario": "Major version upgrade or risky change",
            "recommendation": "Canary",
            "rationale": "Progressive rollout, early issue detection, automated safety"
        },
        {
            "scenario": "Limited infrastructure capacity",
            "recommendation": "Rolling Update or Canary",
            "rationale": "Don't require 2x resources like Blue-Green"
        }
    ]

    for guide in decision_guide:
        add_shape_box(slide, 1.5, y, 3.5, 0.35, guide['scenario'],
                     LIGHT_GRAY, DARK_GRAY, font_size=10, bold=True)
        add_shape_box(slide, 5.1, y, 1.8, 0.35, guide['recommendation'],
                     SUCCESS_GREEN, WHITE, font_size=10, bold=True)
        add_text_box(slide, 7, y + 0.05, 4.5, 0.28, guide['rationale'], font_size=9, color=DARK_GRAY)
        y += 0.45

    # Final recommendation
    y += 0.2
    add_shape_box(slide, 2, y, 9, 0.55,
                 "💡 Recommendation: Start with Rolling for non-prod, " +
                 "graduate to Blue-Green/Canary for production critical services",
                 HARNESS_BLUE, WHITE, font_size=12, bold=False)

def create_best_practices_slide(prs):
    """Slide 24: Best Practices"""
    slide = add_slide_with_title(prs, "Best Practices")

    y = 1.5

    # Best practices categories
    categories = [
        {
            "category": "🔧 VM Deployments",
            "color": HARNESS_BLUE,
            "practices": [
                "Use SSH keys, not passwords - rotate every 90 days",
                "Always create backup before deployment",
                "Test rollback procedures monthly",
                "Implement comprehensive health checks",
                "Monitor deployments for 30 minutes post-deployment"
            ]
        },
        {
            "category": "🐳 ECS Deployments",
            "color": SUCCESS_GREEN,
            "practices": [
                "Version control task definitions in Git",
                "Right-size CPU/memory (monitor and adjust)",
                "Use ALB with proper health check configuration",
                "Enable deployment circuit breaker with auto-rollback",
                "Leverage Container Insights for visibility"
            ]
        },
        {
            "category": "🔒 Security",
            "color": ERROR_RED,
            "practices": [
                "Use secrets management (never hardcode)",
                "Implement RBAC with least privilege",
                "Enable audit logging for all actions",
                "Regular security scans (images and infrastructure)",
                "Keep delegates and tools updated"
            ]
        },
        {
            "category": "📊 Monitoring",
            "color": WARNING_ORANGE,
            "practices": [
                "Set up alerts for critical metrics",
                "Monitor error rates, response times, resource usage",
                "Implement continuous verification with baselines",
                "Create runbooks for common issues",
                "Regular review of deployment metrics"
            ]
        }
    ]

    row = 0
    for i, category in enumerate(categories):
        x = 1 if i % 2 == 0 else 6.5
        y_pos = y + (row * 1.65)

        add_shape_box(slide, x, y_pos, 5, 0.35, category['category'],
                     category['color'], WHITE, font_size=13, bold=True)

        practices_text = "\n".join([f"• {p}" for p in category['practices']])
        add_text_box(slide, x + 0.1, y_pos + 0.4, 4.8, 1.15,
                    practices_text, font_size=10, color=DARK_GRAY)

        if i % 2 == 1:
            row += 1

    # General recommendations
    y += 3.4
    add_text_box(slide, 1, y, 11, 0.3, "General Recommendations:", font_size=16, bold=True, color=DARK_BLUE)

    y += 0.35
    general = [
        "✓ Start small: Begin with dev/QA environments before production",
        "✓ Automate everything: Manual steps are error-prone and slow",
        "✓ Test failure scenarios: Deliberately test rollback and recovery procedures",
        "✓ Document thoroughly: Runbooks, architecture diagrams, troubleshooting guides",
        "✓ Train your team: Regular training sessions and knowledge sharing"
    ]
    add_bullet_list(slide, 1.5, y, 10, 1.2, general, font_size=13)

    # Success metrics
    y += 1.3
    add_shape_box(slide, 2.5, y, 8, 0.5,
                 "📈 Success Metrics: Deployment Frequency ↑ | Deployment Time ↓ | " +
                 "Incident Rate ↓ | MTTR ↓ | Team Confidence ↑",
                 SUCCESS_GREEN, WHITE, font_size=11, bold=False)

def create_implementation_roadmap_slide(prs):
    """Slide 25: Implementation Roadmap"""
    slide = add_slide_with_title(prs, "Implementation Roadmap")

    y = 1.5

    # Timeline phases
    phases = [
        {
            "phase": "Week 1-2: Foundation",
            "icon": "🏗️",
            "tasks": [
                "Set up Harness account",
                "Install delegates",
                "Configure connectors (Git, AWS, Docker)",
                "Create environments (Dev, QA)"
            ],
            "deliverable": "Harness platform ready",
            "color": HARNESS_BLUE
        },
        {
            "phase": "Week 3-4: Development",
            "icon": "⚙️",
            "tasks": [
                "Create service definitions",
                "Build deployment pipelines",
                "Configure deployment strategies",
                "Set up health checks"
            ],
            "deliverable": "Working pipelines in Dev/QA",
            "color": SUCCESS_GREEN
        },
        {
            "phase": "Week 5-6: Testing",
            "icon": "🧪",
            "tasks": [
                "Test deployments in QA",
                "Validate rollback procedures",
                "Performance testing",
                "Security validation"
            ],
            "deliverable": "Tested and validated pipelines",
            "color": WARNING_ORANGE
        },
        {
            "phase": "Week 7-8: Production",
            "icon": "🚀",
            "tasks": [
                "Production environment setup",
                "Final security review",
                "Production deployment",
                "Team training"
            ],
            "deliverable": "Production-ready deployment",
            "color": ERROR_RED
        }
    ]

    for phase in phases:
        # Phase header
        add_shape_box(slide, 1, y, 11, 0.4, f"{phase['icon']} {phase['phase']}",
                     phase['color'], WHITE, font_size=14, bold=True)

        y += 0.45

        # Tasks and deliverable side by side
        tasks_text = "\n".join([f"• {task}" for task in phase['tasks']])
        add_text_box(slide, 1.2, y, 7.5, 0.8, tasks_text, font_size=11, color=DARK_GRAY)

        add_shape_box(slide, 8.8, y + 0.15, 2.8, 0.5, f"✓ {phase['deliverable']}",
                     LIGHT_GRAY, phase['color'], font_size=10, bold=False)

        y += 0.95

    # Resource requirements
    y += 0.1
    add_text_box(slide, 1, y, 11, 0.3, "Resource Requirements:", font_size=16, bold=True, color=DARK_BLUE)

    y += 0.35

    resources = [
        ("DevOps Engineer", "1-2 FTE", "Platform setup, pipeline development"),
        ("Developer", "0.5-1 FTE", "Application configurations"),
        ("QA Engineer", "0.5 FTE", "Testing and validation")
    ]

    x = 1.5
    for role, allocation, responsibility in resources:
        add_shape_box(slide, x, y, 2, 0.3, role, HARNESS_BLUE, WHITE, font_size=11, bold=True)
        add_text_box(slide, x, y + 0.35, 2, 0.15, allocation, font_size=9, color=DARK_GRAY, align=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + 0.52, 2, 0.3, responsibility, font_size=8, color=DARK_GRAY)
        x += 3.2

    # Success criteria
    y += 0.9
    add_shape_box(slide, 2, y, 9, 0.5,
                 "✅ Success Criteria: All pipelines tested | Team trained | " +
                 "Production deployment successful | Monitoring configured",
                 SUCCESS_GREEN, WHITE, font_size=12, bold=False)

def create_benefits_slide(prs):
    """Slide 26: Benefits"""
    slide = add_slide_with_title(prs, "Business Benefits")

    y = 1.5

    # Key benefits with metrics
    benefits = [
        {
            "metric": "70%",
            "label": "Faster\nDeployments",
            "description": "Automated pipelines reduce deployment time from hours to minutes",
            "color": SUCCESS_GREEN
        },
        {
            "metric": "60%",
            "label": "Fewer\nIncidents",
            "description": "Automated verification catches issues before production impact",
            "color": HARNESS_BLUE
        },
        {
            "metric": "Zero",
            "label": "Downtime\nReleases",
            "description": "Blue-Green and Canary strategies ensure continuous availability",
            "color": WARNING_ORANGE
        },
        {
            "metric": "100%",
            "label": "Audit\nCompliance",
            "description": "Complete audit trails for SOC 2, ISO 27001, HIPAA compliance",
            "color": DARK_BLUE
        }
    ]

    x = 1.2
    for benefit in benefits:
        # Metric box
        add_shape_box(slide, x, y, 2.5, 0.7, benefit['metric'],
                     benefit['color'], WHITE, font_size=36, bold=True)

        # Label
        add_text_box(slide, x, y + 0.75, 2.5, 0.4, benefit['label'],
                    font_size=13, bold=True, color=benefit['color'], align=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, x, y + 1.2, 2.5, 0.6, benefit['description'],
                    font_size=10, color=DARK_GRAY)

        x += 2.7

    # Additional benefits
    y += 2
    add_text_box(slide, 1, y, 11, 0.3, "Additional Benefits:", font_size=18, bold=True, color=HARNESS_BLUE)

    y += 0.4

    additional_benefits = [
        {
            "category": "⚡ Speed & Efficiency",
            "benefits": [
                "Deploy to production multiple times per day",
                "Parallel deployments across environments",
                "Automated artifact management"
            ]
        },
        {
            "category": "🎯 Quality & Reliability",
            "benefits": [
                "Consistent deployments across environments",
                "Automated rollback on failures",
                "Comprehensive health checks"
            ]
        },
        {
            "category": "👥 Team Productivity",
            "benefits": [
                "Self-service deployment for developers",
                "Reduced manual intervention",
                "Clear deployment history and audit trails"
            ]
        },
        {
            "category": "💰 Cost Optimization",
            "benefits": [
                "Reduced operations overhead",
                "Fewer production incidents = less firefighting",
                "Efficient resource utilization"
            ]
        }
    ]

    row = 0
    for i, category_info in enumerate(additional_benefits):
        x = 1.2 if i % 2 == 0 else 6.5
        y_pos = y + (row * 1)

        add_text_box(slide, x, y_pos, 5, 0.25, category_info['category'],
                    font_size=14, bold=True, color=HARNESS_BLUE)

        benefits_text = "\n".join([f"• {b}" for b in category_info['benefits']])
        add_text_box(slide, x + 0.2, y_pos + 0.3, 4.8, 0.6,
                    benefits_text, font_size=11, color=DARK_GRAY)

        if i % 2 == 1:
            row += 1

    # ROI statement
    y += 2.2
    add_shape_box(slide, 2, y, 9, 0.6,
                 "💎 Typical ROI: 3-6 months | Reduced MTTR by 80% | " +
                 "Increased deployment frequency by 10x",
                 SUCCESS_GREEN, WHITE, font_size=13, bold=True)

def create_key_takeaways_slide(prs):
    """Slide 27: Key Takeaways"""
    slide = add_slide_with_title(prs, "Key Takeaways")

    y = 1.5

    takeaways = [
        {
            "title": "🏗️ Architecture",
            "points": [
                "Control Plane (SaaS) + Delegates (Customer Network) = Secure & Scalable",
                "No inbound connections required - delegates initiate outbound only",
                "Multi-AZ deployment for high availability"
            ],
            "color": HARNESS_BLUE
        },
        {
            "title": "🚀 VM Deployments",
            "points": [
                "SSH-based deployment to traditional servers",
                "7-phase process: Trigger → Pre-Deploy → Deploy → Post-Deploy → Verify → Complete",
                "Rolling updates with backup and rollback support"
            ],
            "color": DARK_BLUE
        },
        {
            "title": "🐳 ECS Deployments",
            "points": [
                "Container orchestration on AWS Fargate/EC2",
                "Blue-Green for zero-downtime, Canary for risk mitigation",
                "Deep AWS integration (ALB, ECR, CloudWatch)"
            ],
            "color": SUCCESS_GREEN
        },
        {
            "title": "📊 Verification",
            "points": [
                "AI-powered continuous verification compares against baseline",
                "Automatic rollback on anomaly detection",
                "Multi-source monitoring (Prometheus, Datadog, CloudWatch)"
            ],
            "color": WARNING_ORANGE
        },
        {
            "title": "🔒 Security",
            "points": [
                "Network isolation with VPC, security groups, private subnets",
                "Secrets management (AWS Secrets Manager, Vault)",
                "Complete audit trails and compliance (SOC 2, ISO 27001)"
            ],
            "color": ERROR_RED
        }
    ]

    for takeaway in takeaways:
        add_shape_box(slide, 1, y, 11, 0.35, takeaway['title'],
                     takeaway['color'], WHITE, font_size=14, bold=True)

        y += 0.4

        points_text = "\n".join([f"• {point}" for point in takeaway['points']])
        add_text_box(slide, 1.3, y, 10.5, 0.6, points_text, font_size=11, color=DARK_GRAY)

        y += 0.75

    # Final summary
    y += 0.1
    add_shape_box(slide, 1.5, y, 10, 0.7,
                 "✨ Harness CD enables: Faster deployments | Fewer incidents | Zero downtime | " +
                 "Complete visibility | Automated safety | Enterprise security",
                 HARNESS_BLUE, WHITE, font_size=14, bold=True)

def create_next_steps_slide(prs):
    """Slide 28: Next Steps"""
    slide = add_slide_with_title(prs, "Next Steps")

    y = 1.5

    # Immediate actions
    add_shape_box(slide, 1, y, 11, 0.5, "🎯 Immediate Actions (This Week)",
                 HARNESS_BLUE, WHITE, font_size=18, bold=True)

    y += 0.6

    immediate = [
        "✓ Set up Harness account (free trial: https://harness.io/trial)",
        "✓ Identify pilot application/service for initial deployment",
        "✓ Form project team (1-2 DevOps engineers, 1 developer)",
        "✓ Review infrastructure requirements and access permissions"
    ]
    add_bullet_list(slide, 1.5, y, 10, 0.9, immediate, font_size=14)

    # Short-term
    y += 1.1
    add_shape_box(slide, 1, y, 11, 0.5, "📅 Short-term (Next 2 Weeks)",
                 SUCCESS_GREEN, WHITE, font_size=18, bold=True)

    y += 0.6

    short_term = [
        "✓ Install Harness delegate in your environment",
        "✓ Configure connectors (Git, AWS/Cloud provider, Artifact repositories)",
        "✓ Create first service definition and deployment pipeline",
        "✓ Test deployment in Dev environment"
    ]
    add_bullet_list(slide, 1.5, y, 10, 0.9, short_term, font_size=14)

    # Medium-term
    y += 1.1
    add_shape_box(slide, 1, y, 11, 0.5, "🚀 Medium-term (Next 4-6 Weeks)",
                 WARNING_ORANGE, WHITE, font_size=18, bold=True)

    y += 0.6

    medium_term = [
        "✓ Expand to QA and Staging environments",
        "✓ Implement Blue-Green or Canary deployment strategy",
        "✓ Set up continuous verification and monitoring",
        "✓ Plan production deployment"
    ]
    add_bullet_list(slide, 1.5, y, 10, 0.9, medium_term, font_size=14)

    # Contact & Resources
    y += 1.1
    add_text_box(slide, 1, y, 11, 0.3, "📞 Resources & Support:", font_size=16, bold=True, color=DARK_BLUE)

    y += 0.35

    resources = [
        "📚 Documentation: https://developer.harness.io/docs/continuous-delivery",
        "🎓 Harness University: https://university.harness.io",
        "💬 Community: https://community.harness.io",
        "📧 Sales: sales@harness.io | Support: support@harness.io"
    ]
    add_bullet_list(slide, 1.5, y, 10, 0.9, resources, font_size=13)

def create_thank_you_slide(prs):
    """Slide 29: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = HARNESS_BLUE

    # Thank you message
    add_text_box(slide, 1, 2.5, 11.3, 1,
                "Thank You!",
                font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Questions
    add_text_box(slide, 1, 3.7, 11.3, 0.6,
                "Questions?",
                font_size=36, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    # Contact info
    y = 4.8
    contact_info = [
        "📧 Email: sales@harness.io",
        "🌐 Website: https://harness.io",
        "📚 Documentation: https://developer.harness.io",
        "💬 Community: https://community.harness.io"
    ]

    for info in contact_info:
        add_text_box(slide, 2, y, 9.3, 0.35, info,
                    font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
        y += 0.45

# =============================================================================
# MAIN FUNCTION
# =============================================================================

def main():
    """Main function to generate the presentation"""
    print("Creating Professional Harness CD Architecture Presentation...")
    print("This is a customer-sharable, comprehensive presentation")
    print()

    try:
        prs = create_presentation()

        output_file = "/Users/tejasodanapalli/Desktop/Dailywork/Harness_CD_Professional_Architecture.pptx"
        prs.save(output_file)

        print("✅ Presentation created successfully!")
        print(f"✅ File saved: {output_file}")
        print(f"✅ Total slides: 29")
        print()
        print("📊 Slide Breakdown:")
        print("  • Slides 1-5: Introduction & Overview")
        print("  • Slides 6-12: VM Deployment (detailed architecture & flows)")
        print("  • Slides 13-19: ECS Deployment (detailed architecture & strategies)")
        print("  • Slides 20-22: Network, Security & Monitoring")
        print("  • Slides 23-26: Comparison, Best Practices & Benefits")
        print("  • Slides 27-29: Takeaways, Next Steps & Thank You")
        print()
        print("🎨 Features:")
        print("  • Professional color scheme (Harness blue, success green, warning orange)")
        print("  • Clear visual diagrams and architecture flows")
        print("  • Comprehensive coverage of VM and ECS deployments")
        print("  • Customer-ready content (no internal information)")
        print("  • Suitable for executive presentations and technical deep-dives")

    except Exception as e:
        print(f"❌ Error creating presentation: {str(e)}")
        import traceback
        traceback.print_exc()
        raise

if __name__ == "__main__":
    main()
